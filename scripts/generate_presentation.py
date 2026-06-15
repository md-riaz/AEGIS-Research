import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Define colors
    BACKGROUND_COLOR = RGBColor(10, 15, 30)  # Very dark blue
    TEXT_COLOR = RGBColor(240, 240, 240)    # Off-white
    ACCENT_COLOR = RGBColor(0, 180, 255)    # Vibrant blue
    
    def apply_styles(slide, title_text):
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BACKGROUND_COLOR
        
        # Style Title
        title = slide.shapes.title
        title.text = title_text
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = ACCENT_COLOR
            paragraph.alignment = PP_ALIGN.LEFT
            
        # Add Footer
        footer_y = prs.slide_height - Inches(0.5)
        footer_text = "AEGIS Research Project"
        date_text = "May 2026"
        
        # Left footer (Date)
        txDate = slide.shapes.add_textbox(Inches(0.5), footer_y, Inches(3), Inches(0.3))
        txDate.text_frame.text = date_text
        txDate.text_frame.paragraphs[0].font.size = Pt(10)
        txDate.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
        
        # Center footer (Dept)
        txDept = slide.shapes.add_textbox(Inches(3.5), footer_y, Inches(6), Inches(0.3))
        txDept.text_frame.text = footer_text
        txDept.text_frame.paragraphs[0].font.size = Pt(10)
        txDept.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
        txDept.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_bullets(slide, points, level=0):
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = point
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(10)

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_styles(slide, "Natural Language to Dashboard")
    title = slide.shapes.title
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    subtitle = slide.placeholders[1]
    subtitle.text = ("A Safe AI-Assisted Reporting and Widget Generation System\n"
                     "Research Presentation\n"
                     "AEGIS Research Division\n\n"
                     "Presented By: Md. Riaz\n")
    for p in subtitle.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER

    # 2. Outline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "Outline")
    add_bullets(slide, [
        "Introduction", "Problem Statement", "Objectives", 
        "Literature Review / Related Work", "What the research proposes", 
        "Methodology", "System Structure", "E-commerce Example", 
        "Conclusion & Future Work", "References"
    ])

    # 3. Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "Introduction")
    add_bullets(slide, [
        "Most software dashboards are predefined and developer-dependent.",
        "Non-technical users cannot translate insights into SQL steps.",
        "Free-form LLM-generated SQL is flexible but unsafe.",
        "Need for a safe natural-language-driven reporting workflow.",
        "Traditional reporting delays decision-making."
    ])

    # 4. Introduction (Con..)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "Introduction (Con..)")
    add_bullets(slide, [
        "Reporting needs fit into small set of analytics primitives:",
        "• Trend → sales over time",
        "• Rank → top 5 products",
        "• Compare → this month vs last month",
        "• Segment → sales by category",
        "• KPI → total revenue today",
        "Controlled combinations of trusted reporting patterns."
    ])
    # Add Dashboard Preview Image
    img_path = r"C:\Users\mdriaz\.gemini\antigravity\brain\ba6797cc-7cf6-4fea-9f1e-eeda9b965eb5\aegis_dashboard_preview_1777290576414.png"
    if os.path.exists(img_path):
        prs.slides[3].shapes.add_picture(img_path, Inches(6), Inches(2), height=Inches(3.5))

    # 5. Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "Problem Statement")
    add_bullets(slide, [
        "Unsafe & Unscalable: Raw SQL exposes sensitive data.",
        "Unpredictable: AI-generated SQL ignores business definitions.",
        "Ambiguity: Users struggle to express intent clearly.",
        "Lack of Persistence: Focus is on query generation, not reusable widgets.",
        "No complete end-to-end workflow exists."
    ])

    # 6. Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "Objectives")
    add_bullets(slide, [
        "Design a natural language interface for reporting.",
        "Predefined semantic layer of trusted business concepts.",
        "Extract structured intent using LLMs.",
        "Safe query compiler using approved templates.",
        "Generate outputs as reusable widgets.",
        "Evaluate correctness, safety, and usability."
    ])

    # 7. Literature Review
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "Literature Review")
    # Table for Lit Review
    table_slide = prs.slides[6]
    rows, cols = 6, 3
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(4)
    table = table_slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    def set_cell(r, c, text, font_size=12, bold=False):
        cell = table.cell(r, c)
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = TEXT_COLOR
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(20, 30, 50)

    headers = ["Author/Source", "Technique", "Findings"]
    for i, h in enumerate(headers): set_cell(0, i, h, 14, True)
    
    data = [
        ["NaLIR (2014)", "Rule-based NL→SQL", "Heavy reliance on clarification"],
        ["G-SQL (2025)", "Schema-aware NL→SQL", "Poor complex schema generalization"],
        ["TriSQL (2026)", "Multi-stage LLM SQL", "Struggles with complex queries"],
        ["Conversational BI", "Chat-based assistant", "Limited scaling/ambiguity"],
        ["DashBot (2022)", "RL for Dashboards", "Limited real-world use"]
    ]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            set_cell(r+1, c, val)

    # 8. Evolution Path
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "Evolution Path")
    add_bullets(slide, [
        "1. NLIDB (Rule-Based) -> NaLIR (2014)",
        "2. Text-to-SQL (ML/DL) -> NLI4DB (2025)",
        "3. LLM-based SQL -> TriSQL (2026)",
        "4. Conversational BI -> IJERT (2025)"
    ])

    apply_styles(slide, "What This Research Proposes")
    add_bullets(slide, [
        "Constraint-based natural language reporting system.",
        "AI for understanding, System for execution.",
        "Structured analysis plan instead of raw SQL.",
        "Compiled using: Approved metrics, dimensions, and relationships.",
        "Output: Widget-ready results (KPI, Table, Chart)."
    ])

    # 10. Methodology
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "Methodology")
    add_bullets(slide, [
        "Query Interface -> LLM Intent Parser",
        "Semantic Layer -> AI Analysis Planner",
        "Safe Query Compiler -> Visualization Selector",
        "Widget Engine -> Final Dashboard"
    ])
    # Add Architecture Diagram here
    arch_img = r"C:\Users\mdriaz\.gemini\antigravity\brain\ba6797cc-7cf6-4fea-9f1e-eeda9b965eb5\aegis_architecture_diagram_1777290529737.png"
    if os.path.exists(arch_img):
        prs.slides[9].shapes.add_picture(arch_img, Inches(5.5), Inches(2), height=Inches(3.5))

    # 11. System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "System Architecture")
    add_bullets(slide, [
        "Relational Schema (e.g., e-commerce Management).",
        "Semantic Layer: Metrics, Dimensions, Filters.",
        "Pattern Library: Ranking, Trend, Comparison.",
        "LLM API: Intent extraction (Structured JSON).",
        "Safe Query Compiler & Widget Engine.",
        "Evaluation: Accuracy, Safety, Satisfaction."
    ])

    # 12. LEGO vs Clay
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "Structure: LEGO blocks, not Clay")
    add_bullets(slide, [
        "Metrics: What you measure.",
        "Dimensions: How you slice.",
        "Time Rules & Relationships.",
        "User asks: 'Top 5 products by sales last month'",
        "Maps to: total_sales, product, last_month, rank.",
        "Output: Bar Chart."
    ])
    lego_img = r"C:\Users\mdriaz\.gemini\antigravity\brain\ba6797cc-7cf6-4fea-9f1e-eeda9b965eb5\aegis_lego_analogy_1777292513406.png"
    if os.path.exists(lego_img):
        prs.slides[11].shapes.add_picture(lego_img, Inches(6), Inches(2), height=Inches(3.5))

    # 13. AI Freedom
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "AI Freedom in Interpretation")
    add_bullets(slide, [
        "AI CAN: Rephrase intent, infer context, suggest follow-ups.",
        "AI CANNOT: Invent joins, bypass permissions, hit raw tables.",
        "Freedom of thinking, not acting."
    ])

    # 14. Open-ended Growth
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "Designed for Extension")
    add_bullets(slide, [
        "Add new metrics/dimensions easily.",
        "Primitives are rare and stable.",
        "Growth by extension, not reinvention.",
        "Similar to Power BI or Tableau evolution."
    ])

    # 15. E-commerce Example
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "E-commerce Example")
    add_bullets(slide, [
        "Prompt: 'Show top 5 categories by refund rate this month'",
        "Intent: Ranking",
        "Metric: total_revenue (Adjusted for refunds)",
        "Dimension: category_name",
        "Time Rule: current_month",
        "Output: Compiled SQL + Bar Chart + Saved Widget"
    ])

    # 16. Conclusion & Future Work
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "Conclusion & Future Work")
    add_bullets(slide, [
        "Safe architecture for natural-language analytics.",
        "AI assists interpretation without controlling execution.",
        "Future: Implement prototype, RAG integration.",
        "Voice controls & Automated alerts."
    ])
    # Add Performance Chart here
    perf_img = r"C:\Users\mdriaz\.gemini\antigravity\brain\ba6797cc-7cf6-4fea-9f1e-eeda9b965eb5\aegis_performance_chart_1777290651053.png"
    if os.path.exists(perf_img):
        prs.slides[15].shapes.add_picture(perf_img, Inches(6), Inches(2.5), height=Inches(3))

    # 17. References / Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_styles(slide, "References & Questions")
    add_bullets(slide, [
        "[1] Li & Jagadish (2014) NaLIR",
        "[2] Shalaan et al. (2025) G-SQL",
        "[3] Su et al. (2026) TriSQL",
        "[4] Shailesh et al. (2025) Conversational BI",
        "[5] Deng et al. (2023) DashBot",
        "",
        "Thank You. Any Questions?"
    ])

    output_path = "AEGIS_Research_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
