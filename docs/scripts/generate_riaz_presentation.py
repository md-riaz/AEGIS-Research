import os
import re
import copy
import io
import json
from datetime import datetime
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE_FONT = 'Times New Roman'  # matches university template typography throughout
ACCENT_GREEN = RGBColor(0x70, 0xAD, 0x47)  # theme accent6 - matches template's actual brand palette
HEADER_COLOR = RGBColor(0, 51, 102)
BODY_COLOR = RGBColor(51, 51, 51)
SUBDETAIL_COLOR = RGBColor(95, 95, 95)
CODE_COLOR = RGBColor(80, 80, 80)


def _split_label(content):
    """Split 'Label: description' into a bold lead-in and a normal remainder.
    Returns (label_incl_colon, rest_incl_leading_space) or (None, content) if no label pattern exists."""
    if content.endswith(':'):
        return content, ''
    if ': ' in content:
        label, rest = content.split(': ', 1)
        return label + ':', ' ' + rest
    return None, content


def _set_hanging_indent(paragraph, marL_inches, indent_inches):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('marL', str(int(Inches(marL_inches))))
    pPr.set('indent', str(int(Inches(indent_inches))))


def add_fit_picture(slide, image_path, left, top, width, height):
    """Insert an image scaled to fit inside a fixed box without distortion."""
    from PIL import Image
    with Image.open(image_path) as img:
        img_w, img_h = img.size
    box_w = int(width)
    box_h = int(height)
    scale = min(box_w / img_w, box_h / img_h)
    final_w = int(img_w * scale)
    final_h = int(img_h * scale)
    final_left = int(left + (box_w - final_w) / 2)
    final_top = int(top + (box_h - final_h) / 2)
    return slide.shapes.add_picture(image_path, final_left, final_top, final_w, final_h)


# --------------------------------------------------------------------------
# Text measurement. PowerPoint only recomputes its own "shrink text on
# overflow" autofit when a human opens and edits the shape, so a generated
# deck silently spills body text under the footer band. We therefore measure
# the wrapped text ourselves and pick a font size that actually fits.
#
# Pillow ships as a python-pptx dependency, so it is always importable. The
# Linux Liberation fonts are metric-compatible with the Windows fonts the
# template uses, which keeps the measurement valid on either platform.
# --------------------------------------------------------------------------
_FONT_FILES = {
    ('serif', False, False): [r'C:\Windows\Fonts\times.ttf',
                              '/usr/share/fonts/truetype/msttcorefonts/times.ttf',
                              '/usr/share/fonts/truetype/liberation/LiberationSerif-Regular.ttf'],
    ('serif', True, False): [r'C:\Windows\Fonts\timesbd.ttf',
                             '/usr/share/fonts/truetype/msttcorefonts/timesbd.ttf',
                             '/usr/share/fonts/truetype/liberation/LiberationSerif-Bold.ttf'],
    ('serif', False, True): [r'C:\Windows\Fonts\timesi.ttf',
                             '/usr/share/fonts/truetype/msttcorefonts/timesi.ttf',
                             '/usr/share/fonts/truetype/liberation/LiberationSerif-Italic.ttf'],
    ('serif', True, True): [r'C:\Windows\Fonts\timesbi.ttf',
                            '/usr/share/fonts/truetype/msttcorefonts/timesbi.ttf',
                            '/usr/share/fonts/truetype/liberation/LiberationSerif-BoldItalic.ttf'],
    ('mono', False, False): [r'C:\Windows\Fonts\consola.ttf',
                             '/usr/share/fonts/truetype/liberation/LiberationMono-Regular.ttf'],
}
_font_cache = {}

# Collected at build time and reported at the end so a slide can never quietly
# ship with text running under the footer band.
OVERFLOW_WARNINGS = []

# The template's footer band starts here; no content shape may reach into it.
FOOTER_BAND_TOP_INCHES = 6.95


def _load_font(family, bold, italic, size_pt):
    """Return a PIL font for measurement, or None if no usable font file exists."""
    key = (family, bold, italic, round(size_pt, 1))
    if key in _font_cache:
        return _font_cache[key]
    from PIL import ImageFont
    font = None
    for path in _FONT_FILES.get((family, bold, italic), []):
        if os.path.exists(path):
            # PIL sizes in pixels; at 72 dpi one point is one pixel, so points
            # and pixels coincide and the returned widths are already in points.
            font = ImageFont.truetype(path, max(int(round(size_pt)), 1))
            break
    _font_cache[key] = font
    return font


def _text_width_pt(text, family, bold, italic, size_pt):
    font = _load_font(family, bold, italic, size_pt)
    if font is not None:
        return font.getlength(text)
    # Conservative fallback if no font file is available on this machine.
    return len(text) * size_pt * (0.55 if bold else 0.50)


def _wrapped_line_count(runs, first_line_width_pt, body_width_pt):
    """Greedy word wrap across a paragraph's runs; returns the number of lines."""
    words = []  # (text, family, bold, italic, size)
    for text, family, bold, italic, size in runs:
        parts = text.split(' ')
        for i, part in enumerate(parts):
            if part == '' and i > 0:
                continue
            words.append((part, family, bold, italic, size))
    if not words:
        return 1

    lines = 1
    limit = first_line_width_pt
    used = 0.0
    for i, (word, family, bold, italic, size) in enumerate(words):
        piece = word if used == 0 else ' ' + word
        w = _text_width_pt(piece, family, bold, italic, size)
        if used > 0 and used + w > limit:
            lines += 1
            limit = body_width_pt
            used = _text_width_pt(word, family, bold, italic, size)
        else:
            used += w
    return lines


def _build_paragraph_plan(text, font_size, header_color):
    """Parse the structured text block into per-paragraph render instructions.

    Kept separate from rendering so the same plan can be measured at several
    candidate font sizes before anything is written into the slide.
    """
    plan = []
    for raw in text.split('\n'):
        stripped = raw.strip()
        if not stripped:
            continue

        leading_spaces = len(raw) - len(raw.lstrip(' '))
        is_indented = leading_spaces >= 1
        is_bullet = stripped.startswith('- ') and not is_indented
        is_dash_bullet = stripped.startswith('- ') and not is_indented
        is_numbered = bool(re.match(r'^\d+\.\s', stripped)) and not is_indented
        is_sql_line = bool(re.match(r'^(SELECT|FROM|WHERE|GROUP BY|ORDER BY|LIMIT|LEFT JOIN|JOIN|AND)\b', stripped))
        is_code = stripped in ('{', '}') or (is_indented and (stripped.startswith('"') or is_sql_line))

        if is_code:
            plan.append(dict(
                runs=[(stripped, 'Consolas', 'mono', False, False,
                       max(font_size - 3, 11), CODE_COLOR)],
                marL=0.4, indent=0.0, space_before=0, space_after=0, wrap=False))
            continue

        if is_bullet or is_dash_bullet or is_numbered:
            if is_bullet:
                content = stripped[2:].strip()
            elif is_dash_bullet:
                content = stripped[2:].strip()
            else:
                content = stripped
            label, rest = _split_label(content)
            if is_bullet:
                marker = '- '
            elif is_dash_bullet:
                marker = '- '
            else:
                marker = ''
            if label is not None:
                runs = [(marker + label, TEMPLATE_FONT, 'serif', True, False,
                         font_size, header_color)]
                if rest:
                    runs.append((rest, TEMPLATE_FONT, 'serif', False, False,
                                 font_size, BODY_COLOR))
            else:
                runs = [(marker + content, TEMPLATE_FONT, 'serif', False, False,
                         font_size, BODY_COLOR)]
            plan.append(dict(runs=runs, marL=0.28, indent=-0.28,
                             space_before=3, space_after=7, wrap=True))
            continue

        if is_indented:
            plan.append(dict(
                runs=[(stripped, TEMPLATE_FONT, 'serif', False, True,
                       max(font_size - 2, 12), SUBDETAIL_COLOR)],
                marL=0.5, indent=0.0, space_before=0, space_after=8, wrap=True))
            continue

        # Section header / framing statement. A bare "Heading:" gets the full
        # bold treatment; a "Label: narrative sentence" gets a bold lead-in
        # plus a normal-weight continuation so it doesn't read as a wall of bold text.
        label, rest = _split_label(stripped)
        if label is not None:
            runs = [(label, TEMPLATE_FONT, 'serif', True, False,
                     font_size + (3 if not rest else 1), header_color)]
            if rest:
                runs.append((rest, TEMPLATE_FONT, 'serif', False, False,
                             font_size, BODY_COLOR))
        else:
            runs = [(stripped, TEMPLATE_FONT, 'serif', True, False,
                     font_size + 3, header_color)]
        plan.append(dict(runs=runs, marL=0.0, indent=0.0,
                         space_before=12, space_after=5, wrap=True))

    return plan


def _plan_height_pt(plan, width_inches, line_spacing=1.08):
    """Estimated rendered height of a paragraph plan, in points."""
    total = 0.0
    for i, para in enumerate(plan):
        max_size = max(r[5] for r in para['runs'])
        if para['wrap']:
            body_w = (width_inches - para['marL']) * 72.0
            first_w = (width_inches - para['marL'] - para['indent']) * 72.0
            measure_runs = [(r[0], r[2], r[3], r[4], r[5]) for r in para['runs']]
            lines = _wrapped_line_count(measure_runs, first_w, body_w)
        else:
            lines = 1
        # PowerPoint's single-spaced line box is ~1.2x the point size.
        total += lines * max_size * 1.2 * line_spacing
        total += para['space_after']
        if i > 0:
            total += para['space_before']
    return total


def add_bullet_text(slide, text, left, top, width, height, font_size=18, header_color=None,
                    min_font_size=11, autofit=True, line_spacing=1.08):
    """Render a structured block of text with real typographic hierarchy:
    - lines ending in ':' or plain framing statements -> bold section headers
    - lines starting with '- ' or 'N. ' -> hanging-indent bullets, with an optional
      bold 'Label:' lead-in split from the rest of the sentence
    - indented lines -> italic sub-detail/description text
    - lines that are '{', '}', or indented quoted JSON fields -> small monospace code text

    When ``autofit`` is on, the requested ``font_size`` is the maximum: the text
    is measured and stepped down (never below ``min_font_size``) until it fits
    inside ``height``, so no slide can spill body text under the footer band.
    """
    if header_color is None:
        header_color = HEADER_COLOR

    width_inches = width / 914400
    available_pt = (height / 914400) * 72.0

    effective_size = font_size
    if autofit:
        candidate = font_size
        while candidate > min_font_size:
            plan = _build_paragraph_plan(text, candidate, header_color)
            if _plan_height_pt(plan, width_inches, line_spacing) <= available_pt:
                break
            candidate -= 0.5
        effective_size = max(candidate, min_font_size)

    plan = _build_paragraph_plan(text, effective_size, header_color)
    overflow = _plan_height_pt(plan, width_inches, line_spacing) - available_pt
    if autofit and overflow > 0:
        # Autofit bottomed out at min_font_size and the text still does not fit;
        # the slide needs trimming rather than a smaller font.
        OVERFLOW_WARNINGS.append(
            "%.0fpt of text overflows a %.2f\" box at the %spt floor: %r"
            % (overflow, height / 914400, min_font_size, text.strip()[:60]))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_right = 0

    for i, para in enumerate(plan):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        for run_text, font_name, _family, bold, italic, size, color in para['runs']:
            run = p.add_run()
            run.text = run_text
            run.font.name = font_name
            run.font.bold = bold
            run.font.italic = italic
            run.font.size = Pt(size)
            run.font.color.rgb = color
        _set_hanging_indent(p, para['marL'], para['indent'])
        p.space_before = Pt(para['space_before'])
        p.space_after = Pt(para['space_after'])

    return txBox


def add_problem_statement_text(slide, left, top, width, height, font_size=16):
    """Render the Problem Statement slide without blue body-text headers."""
    lines = [
        ("Current Generative NL2SQL approaches have 3 structural vulnerability classes:", True),
        ("1. Structural Injection Risk", True),
        ("Adversarial prompt manipulation can bypass model instructions, causing neural models to output data-modifying DML/DDL statements (DROP, DELETE, UPDATE).", False),
        ("2. Unbounded Schema Hallucination", True),
        ("Probabilistic token generation leads to hallucinated table joins, non-existent entity relations, and invalid column attributes.", False),
        ("3. Access Control & Context Bypass", True),
        ("Direct query generation bypasses application-level multi-tenant boundaries and row-level security scopes.", False),
    ]
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_right = 0

    for i, (text, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.08
        p.space_before = Pt(8 if i in (1, 3, 5) else 0)
        p.space_after = Pt(5 if bold else 9)
        run = p.add_run()
        run.text = text
        run.font.name = TEMPLATE_FONT
        run.font.bold = bold
        run.font.italic = False
        run.font.size = Pt(font_size + (1 if i == 0 else 0))
        run.font.color.rgb = BODY_COLOR
    return txBox


def style_table(table, header_rows=1, zebra_color=RGBColor(0xF2, 0xF4, 0xF8),
                margin_left=0.1, margin_right=0.1, margin_top=0.05, margin_bottom=0.05):
    """Vertically center cell text, add breathing-room margins, and zebra-stripe body rows."""
    for r_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(margin_left)
            cell.margin_right = Inches(margin_right)
            cell.margin_top = Inches(margin_top)
            cell.margin_bottom = Inches(margin_bottom)
            if r_idx >= header_rows and (r_idx - header_rows) % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = zebra_color

def create_presentation():
    """Build the final-defense deck using the university template style."""
    script_dir = os.path.dirname(os.path.abspath(__file__))
    repo_root = os.path.dirname(os.path.dirname(script_dir))
    output_path = os.path.join(script_dir, 'Md_Riaz_Final_Defense_0322310105101024.pptx')
    figures_dir = os.path.join(script_dir, 'thesis_book_generator', 'figures')
    presentation_assets_dir = os.path.join(script_dir, 'presentation_assets')
    template_path = os.path.join(
        repo_root, 'other person thesis', 'Md.Mominur Rahaman spring2022 Batch 14th Id 0322210105101511.pptx')

    def load_result_json(filename):
        """Read a committed evaluation artifact, or return None if it is absent.

        The results slides render from these files rather than from figures
        typed into this script, so a re-run followed by a rebuild updates the
        deck automatically and the two cannot disagree. Returning None on a
        missing file matters as much as the happy path: the slide then prints
        "not measured" instead of leaving a previous run's number on screen,
        which is the failure this project keeps correcting elsewhere.
        """
        path = os.path.join(repo_root, 'evaluation_dataset', filename)
        try:
            with open(path, encoding='utf-8-sig') as fh:
                return json.load(fh)
        except (OSError, ValueError):
            return None

    prs = pptx.Presentation(template_path)
    prs.core_properties.title = 'AEGIS: A Constraint-Based Architecture for Safe LLM-Assisted Natural Language Analytics'
    prs.core_properties.subject = 'Final Defense Presentation'
    prs.core_properties.author = 'Md. Riaz'
    orig_s1, orig_s2 = list(prs.slides)[:2]

    logo_blob_s1 = [s for s in orig_s1.shapes if s.name == 'Picture 9'][0].image.blob
    logo_blob_s2 = [s for s in orig_s2.shapes if s.name == 'Picture 9'][0].image.blob
    pic_s1 = [s for s in orig_s1.shapes if s.name == 'Picture 9'][0]
    pic_s2 = [s for s in orig_s2.shapes if s.name == 'Picture 9'][0]
    footer_placeholders = [sh for sh in orig_s2.shapes
                           if sh.name in ('Date Placeholder 3', 'Footer Placeholder 4', 'Slide Number Placeholder 5')]
    dept_textbox = [sh for sh in orig_s1.shapes if sh.name == 'TextBox 14'][0]
    # Shown in the footer of every content slide. Set this to the defense date
    # before the deck is presented — a stale date is the first thing a committee
    # notices, and it appears on all twenty-odd slides.
    footer_date = "Tuesday, August 18, 2026"

    for s in list(prs.slides._sldIdLst):
        prs.part.drop_rel(s.rId)
    prs.slides._sldIdLst.clear()

    primary_color = RGBColor(0, 51, 102)
    light_blue = RGBColor(0xE9, 0xF1, 0xFA)
    soft_green = RGBColor(0xE2, 0xF0, 0xD9)

    def apply_title_slide_branding(s):
        for sh in orig_s1.shapes:
            if sh.name in ['Rectangle 2', 'Rectangle 6']:
                s.shapes._spTree.insert(2, copy.deepcopy(sh._element))
        s.shapes.add_picture(io.BytesIO(logo_blob_s1), pic_s1.left, pic_s1.top, pic_s1.width, pic_s1.height)
        dept = s.shapes.add_textbox(dept_textbox.left, Inches(5.92), dept_textbox.width, Inches(0.58))
        tf = dept.text_frame
        tf.word_wrap = True
        lines = [
            ("Department of Computer Science & Engineering", 15),
            ("Pundra University of Science & Technology, Bogura - 5800", 12),
        ]
        for i, (line_text, size) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line_text
            p.alignment = PP_ALIGN.CENTER
            p.font.name = TEMPLATE_FONT
            p.font.size = Pt(size)


    def apply_content_slide_branding(s, slide_number):
        for sh in orig_s2.shapes:
            if sh.name in ['Rectangle 2', 'Rectangle 6', 'Rectangle 10']:
                s.shapes._spTree.insert(2, copy.deepcopy(sh._element))
        s.shapes.add_picture(io.BytesIO(logo_blob_s2), pic_s2.left, pic_s2.top, pic_s2.width, pic_s2.height)
        for sh in footer_placeholders:
            s.shapes._spTree.append(copy.deepcopy(sh._element))
        for ph in s.placeholders:
            ph_type = str(ph.placeholder_format.type)
            if ph_type.startswith('DATE'):
                ph.text_frame.text = footer_date
            elif ph_type.startswith('FOOTER'):
                ph.text_frame.text = "Department of Computer Science & Engineering, PUB"
            elif ph_type.startswith('SLIDE_NUMBER'):
                ph.text_frame.text = str(slide_number)
            for p in ph.text_frame.paragraphs:
                p.font.name = TEMPLATE_FONT

    def add_title_slide(notes=""):
        s = prs.slides.add_slide(prs.slide_layouts[0])
        apply_title_slide_branding(s)
        eyebrow = s.shapes.add_textbox(Inches(1.6666), Inches(1.8285), Inches(10.0), Inches(0.4429))
        ep = eyebrow.text_frame.paragraphs[0]
        ep.text = "Final Defense Presentation on"
        ep.alignment = PP_ALIGN.CENTER
        ep.font.size = Pt(22)
        ep.font.name = TEMPLATE_FONT

        title_shape = s.placeholders[0]
        title_shape.left = Inches(1.6666)
        title_shape.top = Inches(2.55)
        title_shape.width = Inches(10.0)
        title_shape.height = Inches(1.05)
        title_shape.text = "AEGIS: A Constraint-Based Architecture for Safe\nLLM-Assisted Natural Language Analytics"
        for p in title_shape.text_frame.paragraphs:
            p.font.color.rgb = primary_color
            p.font.bold = True
            p.font.size = Pt(31)
            p.font.name = TEMPLATE_FONT
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.12

        def add_label_value_block(label_text, value_lines, left_in, label_top_in, value_top_in, width_in):
            label_box = s.shapes.add_textbox(Inches(left_in), Inches(label_top_in), Inches(width_in), Inches(0.4))
            lp = label_box.text_frame.paragraphs[0]
            lp.text = label_text
            lp.font.size = Pt(16)
            lp.font.name = TEMPLATE_FONT
            lp.font.italic = True

            value_box = s.shapes.add_textbox(Inches(left_in), Inches(value_top_in), Inches(width_in), Inches(1.25))
            vtf = value_box.text_frame
            vtf.word_wrap = True
            for i, (line_text, size, bold) in enumerate(value_lines):
                p = vtf.paragraphs[0] if i == 0 else vtf.add_paragraph()
                p.text = line_text
                p.font.size = Pt(size)
                p.font.name = TEMPLATE_FONT
                p.font.bold = bold
                p.line_spacing = 1.12

        add_label_value_block(
            "Presented By", [
                ("Md. Riaz", 20, True),
                ("Program: B.Sc. in CSE (Diploma)", 16, False),
                ("ID: 0322310105101024", 16, False),
                ("Batch: 16th", 16, False),
                ("8th Semester / 4th Year", 16, False),
                ("Session: Spring - 2023", 16, False),
            ], 1.67, 4.08, 4.42, 4.25)
        add_label_value_block(
            "Supervised By", [
                ("Mst. Sahela Rahman", 20, True),
                ("Lecturer", 16, False),
                ("Dept. of CSE, PUB", 16, False),
            ], 7.95, 4.12, 4.52, 4.25)
        if notes:
            s.notes_slide.notes_text_frame.text = notes
        return s

    content_slide_count = [1]

    def add_content_slide(title_text, notes=""):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        content_slide_count[0] += 1
        apply_content_slide_branding(s, content_slide_count[0])
        if len(s.placeholders) > 0:
            title_shape = s.placeholders[0]
            title_shape.left = Inches(1.3)
            title_shape.top = Inches(0.4)
            title_shape.width = Inches(10.5)
            title_shape.height = Inches(0.8)
            title_shape.text = title_text
            for p in title_shape.text_frame.paragraphs:
                p.font.color.rgb = primary_color
                p.font.bold = True
                p.font.size = Pt(24)
                p.font.name = TEMPLATE_FONT
        if notes:
            s.notes_slide.notes_text_frame.text = notes
        return s

    def set_cell(cell, text, size=12, bold=False, color=BODY_COLOR, align=PP_ALIGN.LEFT, fill=None):
        cell.text = text
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        if fill is not None:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
        for p in cell.text_frame.paragraphs:
            p.font.name = TEMPLATE_FONT
            p.font.size = Pt(size)
            p.font.bold = bold
            p.font.color.rgb = color
            p.alignment = align
            p.line_spacing = 1.0

    def add_header_row(table, headers, size=12):
        for i, header in enumerate(headers):
            set_cell(table.cell(0, i), header, size=size, bold=True,
                     color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, fill=primary_color)

    def add_flat_box(slide, text, left, top, width, height, fill, border=primary_color, font_size=15):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.06)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = TEMPLATE_FONT
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = primary_color
        p.alignment = PP_ALIGN.CENTER
        return shape

    def add_sql_box(slide, sql, left, top, width, height, font_size=11):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF7, 0xF7, 0xF7)
        box.line.color.rgb = RGBColor(0xB7, 0xB7, 0xB7)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.04)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.08)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        for i, line in enumerate(sql.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = 'Consolas'
            p.font.size = Pt(font_size)
            p.font.color.rgb = CODE_COLOR
            p.space_after = Pt(1)
            p.alignment = PP_ALIGN.LEFT
        return box

    def add_chart_title(slide, title, left, top, width, height=0.28):
        tx = slide.shapes.add_textbox(left, top, width, Inches(height))
        p = tx.text_frame.paragraphs[0]
        p.text = title
        p.font.name = TEMPLATE_FONT
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = primary_color
        p.alignment = PP_ALIGN.CENTER
        return tx

    def add_bar_chart(slide, title, categories, values, left, top, width, height,
                      axis_max=None, orientation="vertical"):
        add_chart_title(slide, title, left, top, width)
        chart_top = top + Inches(0.42)
        chart_height = height - Inches(0.42)
        axis_max = axis_max or max(values)
        plot_left = left + Inches(0.35 if orientation == "vertical" else 1.25)
        plot_top = chart_top + Inches(0.1)
        plot_width = width - Inches(0.65 if orientation == "vertical" else 1.55)
        plot_height = chart_height - Inches(0.65 if orientation == "vertical" else 0.35)
        axis = slide.shapes.add_connector(1, plot_left, plot_top + plot_height, plot_left + plot_width, plot_top + plot_height)
        axis.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
        axis.line.width = Pt(1)
        if orientation == "vertical":
            bar_gap = Inches(0.12)
            bar_slot = plot_width / len(values)
            for i, (cat, value) in enumerate(zip(categories, values)):
                x = plot_left + i * bar_slot + bar_gap
                bar_w = bar_slot - 2 * bar_gap
                bar_h = plot_height * (value / axis_max)
                y = plot_top + plot_height - bar_h
                rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, bar_w, bar_h)
                rect.fill.solid()
                rect.fill.fore_color.rgb = primary_color
                rect.line.color.rgb = primary_color
                label = slide.shapes.add_textbox(x, y - Inches(0.22), bar_w, Inches(0.2))
                p = label.text_frame.paragraphs[0]
                p.text = f"{value}"
                p.font.name = TEMPLATE_FONT
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = BODY_COLOR
                p.alignment = PP_ALIGN.CENTER
                cat_box = slide.shapes.add_textbox(x - Inches(0.08), plot_top + plot_height + Inches(0.06), bar_w + Inches(0.16), Inches(0.42))
                p = cat_box.text_frame.paragraphs[0]
                p.text = cat.replace("\n", "\n")
                p.font.name = TEMPLATE_FONT
                p.font.size = Pt(10)
                p.font.color.rgb = BODY_COLOR
                p.alignment = PP_ALIGN.CENTER
        else:
            bar_gap = Inches(0.12)
            bar_slot = plot_height / len(values)
            for i, (cat, value) in enumerate(zip(categories, values)):
                y = plot_top + i * bar_slot + bar_gap
                bar_h = bar_slot - 2 * bar_gap
                bar_w = plot_width * (value / axis_max)
                cat_box = slide.shapes.add_textbox(left, y - Inches(0.02), Inches(1.15), bar_h + Inches(0.04))
                p = cat_box.text_frame.paragraphs[0]
                p.text = cat
                p.font.name = TEMPLATE_FONT
                p.font.size = Pt(10)
                p.font.color.rgb = BODY_COLOR
                p.alignment = PP_ALIGN.RIGHT
                rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, plot_left, y, bar_w, bar_h)
                rect.fill.solid()
                rect.fill.fore_color.rgb = primary_color
                rect.line.color.rgb = primary_color
                label = slide.shapes.add_textbox(plot_left + bar_w + Inches(0.06), y - Inches(0.02), Inches(0.65), bar_h + Inches(0.04))
                p = label.text_frame.paragraphs[0]
                p.text = f"{value}"
                p.font.name = TEMPLATE_FONT
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = BODY_COLOR

    def add_stage_tag(slide, text):
        """Mark which pipeline stage a technical slide is describing.

        The committee includes non-specialists. Without a fixed marker they have
        no way to tell, eight slides in, whether they are looking at a new stage
        or a closer view of the previous one.
        """
        box = slide.shapes.add_textbox(Inches(10.05), Inches(0.52), Inches(2.3), Inches(0.34))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = TEMPLATE_FONT
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = SUBDETAIL_COLOR
        return box

    def add_plain_line(slide, text, top_inches=6.28, font_size=13):
        """One jargon-free sentence, in the same place on every technical slide.

        Slide 13 already carried a line like this and it was the clearest thing
        on the slide; the value comes from it being in the same position, in the
        same words-per-line budget, every time.
        """
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(top_inches), Inches(11.5), Inches(0.52))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(0xF3, 0xF6, 0xFB)
        bar.line.color.rgb = RGBColor(0xC9, 0xD6, 0xE8)
        bar.line.width = Pt(0.75)
        tf = bar.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.14)
        tf.margin_right = Inches(0.14)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "In plain terms:  "
        run.font.bold = True
        run2 = p.add_run()
        run2.text = text
        for r in (run, run2):
            r.font.name = TEMPLATE_FONT
            r.font.size = Pt(font_size)
            r.font.color.rgb = BODY_COLOR
        p.alignment = PP_ALIGN.LEFT
        return bar

    # 1
    add_title_slide(notes="Good morning/afternoon respected chairperson and committee members. I am presenting the final defense of AEGIS, a constraint-based architecture for safe LLM-assisted natural language analytics.")

    # 2
    s = add_content_slide("Outline")
    add_bullet_text(s, "1. Introduction and Problem Statement\n2. The Idea in One Picture\n3. Literature Review and Research Gap\n4. Objectives and Methodology\n5. AEGIS Architecture (Stages 1-5)\n6. Semantic Layer and Compiler\n7. Query-to-SQL Walkthrough, Including Joins\n8. Prototype Implementation\n9. Evaluation: 500 Questions and 20 Admin Reports\n10. Limitations, Conclusion, and Future Work\n11. Appendix: Schema, Semantic Layer, Compiler Code", Inches(1.55), Inches(1.75), Inches(10.3), Inches(4.8), font_size=18)

    # 3
    s = add_content_slide("Introduction")
    add_bullet_text(s, "Motivation:\n- E-commerce systems store rich operational data, but built-in dashboards answer only a fixed set of questions.\n- When users need new analytical combinations, the request usually depends on a developer or BI person to translate it into report logic.\n- Natural language analytics can reduce that delay for managers, sales teams, and administrators.\n\nCentral Tension:\n- Direct LLM-to-SQL systems are flexible, but the model writes executable database code.\n- A safe analytics system should understand the question without giving the model authority to create arbitrary SQL.\n\nAEGIS Direction:\n- Use the LLM for intent extraction only.\n- Use deterministic semantic mapping and SQL compilation for execution.", Inches(1.2), Inches(1.62), Inches(11.0), Inches(5.18), font_size=16)

    # 4
    s = add_content_slide("Problem Statement")
    add_problem_statement_text(s, Inches(1.2), Inches(1.92), Inches(11.0), Inches(4.38), font_size=16)

    # 5 — the whole thesis in one picture, before any technical slide.
    # Placed here deliberately: everything after this can be checked against it,
    # so a non-specialist listener has somewhere to stand for the next ten
    # slides instead of tracking eight unfamiliar stage names in the abstract.
    s = add_content_slide("The Idea in One Picture")
    add_bullet_text(
        s,
        "In a restaurant, the waiter writes down what you want — the dish, how many plates. "
        "The waiter never walks into the kitchen and cooks it.",
        Inches(1.1), Inches(1.5), Inches(11.1), Inches(0.8), font_size=17)
    add_flat_box(s, "Waiter\ntakes the order", Inches(1.25), Inches(2.5), Inches(3.0), Inches(1.0), light_blue, primary_color, font_size=15)
    add_flat_box(s, "Order form\nfixed fields only", Inches(5.0), Inches(2.5), Inches(3.0), Inches(1.0), soft_green, primary_color, font_size=15)
    add_flat_box(s, "Kitchen\ncooks the dish", Inches(8.75), Inches(2.5), Inches(3.0), Inches(1.0), light_blue, primary_color, font_size=15)
    for x1, x2 in [(4.25, 5.0), (8.0, 8.75)]:
        line = s.shapes.add_connector(1, Inches(x1), Inches(3.0), Inches(x2), Inches(3.0))
        line.line.color.rgb = primary_color
        line.line.width = Pt(2)
    tbl = s.shapes.add_table(4, 3, Inches(1.25), Inches(3.85), Inches(10.5), Inches(2.05)).table
    for i, w in enumerate([3.0, 3.75, 3.75]):
        tbl.columns[i].width = Inches(w)
    add_header_row(tbl, ["In the restaurant", "In AEGIS", "What it cannot do"], size=12.5)
    rows = [
        ["The waiter", "The language model", "Cannot write or run SQL"],
        ["The order form", "The Intent Object", "Cannot hold a table or column name of its own"],
        ["The kitchen", "The SQL compiler", "Cannot accept anything outside the approved list"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            set_cell(tbl.cell(r, c), value, size=12.5, bold=(c == 0),
                     fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl, margin_left=0.07, margin_right=0.07, margin_top=0.03, margin_bottom=0.03)
    add_plain_line(s, "The model chooses names from an approved list. It never writes the database query itself.")

    # 5
    s = add_content_slide("Literature Review")
    tbl = s.shapes.add_table(7, 4, Inches(0.78), Inches(1.55), Inches(11.82), Inches(4.78)).table
    widths = [2.05, 2.65, 3.2, 3.85]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    add_header_row(tbl, ["Work", "Main focus", "Control approach", "Gap for AEGIS"], size=12)
    literature_rows = [
        ["NaLIR [1]", "Interactive NLIDB", "Parse and refine intent with grammar and DB mapping", "Limited dashboard and modern LLM risk handling"],
        ["nl4dv [2]", "NL to visualization specs", "Extract analytic attributes; not SQL-centered", "Visualization focus, not safe database execution"],
        ["DashBot [3]", "Dashboard generation", "Insight-driven dashboard selection", "Does not solve arbitrary SQL authority risk"],
        ["PICARD [4]", "Constrained decoding", "Parser-level SQL validity during token generation", "Model still generates SQL text"],
        ["G-SQL [5] / TriSQL [6]", "Robust Text-to-SQL", "Schema-aware generation with rules, repair, and refinement", "Safety depends on controlling generated SQL"],
        ["AEGIS", "Safe NL analytics", "Intent extraction only; SQL compiled from semantic layer", "Trades open SQL for auditable analytics intent validation"],
    ]
    for r, row in enumerate(literature_rows, start=1):
        for c, text in enumerate(row):
            set_cell(tbl.cell(r, c), text, size=12, bold=(c == 0), fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl, margin_left=0.04, margin_right=0.04, margin_top=0.015, margin_bottom=0.015)

    # 6
    s = add_content_slide("Research Gap")
    add_bullet_text(s, "Observed Gap in Existing Systems:\n- Many systems improve SQL generation accuracy, but still let executable SQL be authored by the model.\n- Visualization systems help users express analysis, but do not fully address database execution safety.\n- Parser and repair methods can reject invalid syntax, yet syntactically valid SQL can still be unsafe or semantically wrong.\n\nAEGIS Research Gap:\n- A practical architecture is needed where natural language is accepted, but executable SQL is produced only from an explicit semantic layer and deterministic compiler.", Inches(1.25), Inches(1.75), Inches(10.9), Inches(4.75), font_size=18)

    # 7
    s = add_content_slide("Objectives and Contributions")
    tbl = s.shapes.add_table(5, 2, Inches(0.9), Inches(1.45), Inches(11.55), Inches(5.0)).table
    tbl.columns[0].width = Inches(3.3)
    tbl.columns[1].width = Inches(8.25)
    add_header_row(tbl, ["Objective", "Contribution"], size=13)
    rows = [
        ["Reduce execution risk", "Separate natural-language understanding from SQL creation."],
        ["Make analytics auditable", "Define metrics, dimensions, filters, join paths, and output shapes in a semantic layer."],
        ["Support real e-commerce reports", "Implement AEGIS over a seeded nopCommerce [7] MySQL database."],
        ["Evaluate against an independent oracle", "Measure breadth on 500 fixed natural-language questions, and fidelity against nopCommerce's own twenty admin reports."],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            set_cell(tbl.cell(r, c), text, size=13, bold=(c == 0), fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl)

    # 8
    s = add_content_slide("Research Methodology")
    fig = os.path.join(presentation_assets_dir, 'thesis-figure-01-dsr-workflow.png')
    if os.path.exists(fig):
        add_fit_picture(s, fig, Inches(1.0), Inches(1.72), Inches(11.35), Inches(4.95))

    # 9
    s = add_content_slide("Proposed AEGIS Architecture")
    fig = os.path.join(figures_dir, 'figure-03-architecture-pipeline.png')
    if os.path.exists(fig):
        add_fit_picture(s, fig, Inches(0.68), Inches(1.2), Inches(12.0), Inches(5.55))

    # 10
    s = add_content_slide("Semantic Layer Implementation")
    fig = os.path.join(figures_dir, 'figure-04-semantic-layer-modularity.png')
    if os.path.exists(fig):
        add_fit_picture(s, fig, Inches(0.7), Inches(1.35), Inches(5.8), Inches(4.75))
    add_bullet_text(s, "Implemented per target system:\n- Business vocabulary: synonyms for metrics, dimensions, filters, and time periods.\n- Metric definitions: SQL expressions such as order count, revenue, average order value, and customer count.\n- Dimension definitions: status, customer, product, manufacturer, category, date, and geography.\n- Join policy: allowed tables and join paths from the nopCommerce [7] schema.\n- Output contracts: table, matrix, KPI, and chart-ready shapes.", Inches(6.75), Inches(1.78), Inches(5.5), Inches(4.55), font_size=14.5)

    # 11
    s = add_content_slide("Intent Extraction Boundary")
    add_stage_tag(s, "Stage 1 of 5")
    add_flat_box(s, "1. User asks a business question", Inches(0.72), Inches(1.46), Inches(3.15), Inches(0.62), light_blue, primary_color, font_size=13)
    add_bullet_text(
        s,
        "Example:\nShow order average by status",
        Inches(0.9), Inches(2.24), Inches(2.9), Inches(0.86),
        font_size=15,
        line_spacing=1.0,
    )
    add_flat_box(s, "2. Prompt exposes approved vocabulary", Inches(4.25), Inches(1.46), Inches(3.65), Inches(0.62), light_blue, primary_color, font_size=12.6)
    add_bullet_text(
        s,
        "Vocabulary excerpt:\n- Metrics: revenue, order_count, avg_order_value\n- Dimensions: order_status, category, country\n- Output: kpi, table, matrix, chart",
        Inches(4.32), Inches(2.2), Inches(3.45), Inches(1.42),
        font_size=10.7,
        line_spacing=0.9,
        min_font_size=9.5,
    )
    add_flat_box(s, "3. LLM returns an Intent Object", Inches(8.62), Inches(1.46), Inches(3.5), Inches(0.62), soft_green, primary_color, font_size=12.8)
    add_sql_box(
        s,
        '{\n  "intent_class": "segment",\n  "metric_term": "avg_order_value",\n  "dimension_term": "order_status",\n  "output_shape": "table"\n}',
        Inches(8.72), Inches(2.2), Inches(3.25), Inches(1.68),
        font_size=9.7,
    )
    for x1, y1, x2, y2 in [
        (3.87, 1.77, 4.25, 1.77),
        (7.9, 1.77, 8.62, 1.77),
    ]:
        line = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = primary_color
        line.line.width = Pt(2)
    tbl = s.shapes.add_table(3, 3, Inches(0.88), Inches(4.68), Inches(11.25), Inches(1.42)).table
    tbl.columns[0].width = Inches(2.45)
    tbl.columns[1].width = Inches(4.1)
    tbl.columns[2].width = Inches(4.7)
    add_header_row(tbl, ["Stage", "Allowed", "Not allowed"], size=11.5)
    boundary_rows = [
        ["LLM", "Choose approved ids from vocabulary", "Raw SQL, joins, table names, writes"],
        ["AEGIS", "Bind ids to semantic definitions", "Unknown ids, silent fallback, unsafe execution"],
    ]
    for r, row in enumerate(boundary_rows, start=1):
        for c, value in enumerate(row):
            set_cell(tbl.cell(r, c), value, size=10.6, bold=(c == 0), align=PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT, fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl, margin_left=0.04, margin_right=0.04, margin_top=0.015, margin_bottom=0.015)
    add_plain_line(s, "The model picks names from a list we wrote. It never sees a table name, and never writes SQL.")

    # 12 — the Intent Object's actual shape, and how a pattern chooses a
    # compilation path. This is the slide that answers "is this really typed,
    # or is it a JSON blob you hope is well-formed?".
    s = add_content_slide("The Intent Object, and What It Selects")
    add_stage_tag(s, "Stage 1 to 2")
    add_flat_box(s, "IntentObject  (aegis/server/models.py)", Inches(0.72), Inches(1.4), Inches(5.5), Inches(0.5), soft_green, primary_color, font_size=12.5)
    add_sql_box(
        s,
        'class IntentObject(BaseModel):\n'
        '    intent_class:    IntentClass      # required, 11-value enum\n'
        '    metric_term:     Optional[str]    # an approved metric id\n'
        '    dimension_term:  Optional[str]    # an approved dimension id\n'
        '    time_term:       Optional[str]    # a phrase, not a date\n'
        '    filters:         List[Filter]     # field / operator / value\n'
        '    limit:           Optional[int]\n'
        '    confidence:      Confidence = LOW # absence is not confidence\n'
        '    unmapped_terms:  List[str]        # what it could not account for\n'
        '\n'
        '# There is no sql field. The type cannot express one.',
        Inches(0.78), Inches(2.0), Inches(5.4), Inches(2.75),
        font_size=9.0,
    )
    add_flat_box(s, "intent_class selects the compilation path", Inches(6.55), Inches(1.4), Inches(5.7), Inches(0.5), light_blue, primary_color, font_size=12.5)
    tbl = s.shapes.add_table(6, 2, Inches(6.6), Inches(2.0), Inches(5.65), Inches(2.75)).table
    tbl.columns[0].width = Inches(2.55)
    tbl.columns[1].width = Inches(3.1)
    add_header_row(tbl, ["intent_class", "Compiler path taken"], size=11)
    rows = [
        ["kpi, summary", "aggregate, no LIMIT"],
        ["segment, ranking,\ncomparison, cohort", "aggregate + GROUP BY + ORDER BY"],
        ["trend", "date bucket; zero-filled when daily"],
        ["tabular, exception", "row listing, no aggregation"],
        ["funnel, correlate", "aggregate over the stated relation"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            set_cell(tbl.cell(r, c), value, size=10, bold=(c == 0),
                     fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl, margin_left=0.05, margin_right=0.05, margin_top=0.02, margin_bottom=0.02)
    add_bullet_text(
        s,
        "An id the semantic layer does not define is rejected here, before any SQL exists. "
        "There is no path from an unknown id to a query.",
        Inches(0.85), Inches(4.95), Inches(11.4), Inches(0.6), font_size=13)
    add_plain_line(s, "The answer form is chosen from eleven fixed shapes, not invented per question.")

    # 13
    s = add_content_slide("Semantic Mapping Under the Hood")
    add_stage_tag(s, "Stage 2 of 5")
    add_flat_box(s, "Intent Object\nmetric_term=avg_order_value\ndimension_term=order_status", Inches(0.82), Inches(1.48), Inches(3.1), Inches(0.92), light_blue, font_size=12)
    add_flat_box(s, "Lookup in semantic layer\nMETRICS + DIMENSIONS", Inches(5.0), Inches(1.48), Inches(3.1), Inches(0.92), soft_green, font_size=12)
    add_flat_box(s, "Analysis Plan\npattern=segment\nshape=table", Inches(9.12), Inches(1.48), Inches(2.55), Inches(0.92), light_blue, font_size=12)
    for x1, x2 in [(3.92, 5.0), (8.1, 9.12)]:
        line = s.shapes.add_connector(1, Inches(x1), Inches(1.94), Inches(x2), Inches(1.94))
        line.line.color.rgb = primary_color
        line.line.width = Pt(2)
    tbl = s.shapes.add_table(4, 5, Inches(0.75), Inches(2.85), Inches(11.85), Inches(2.35)).table
    widths = [1.65, 2.0, 3.1, 2.6, 2.5]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    add_header_row(tbl, ["Slot", "Resolved id", "Business meaning", "SQL expression", "Binding"], size=10.5)
    rows = [
        ["Metric", "avg_order_value", "Average revenue per order", "AVG(COALESCE(o.OrderTotal, 0))", "Order as o"],
        ["Dimension", "order_status", "Readable order status label", "CASE o.OrderStatusId ... END", "Order as o"],
        ["Rule", "mandatory", "Ignore soft-deleted orders", "o.Deleted = 0", "WHERE"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            set_cell(tbl.cell(r, c), value, size=9.8, bold=(c == 0), align=PP_ALIGN.CENTER if c in (0, 1) else PP_ALIGN.LEFT, fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl, margin_left=0.035, margin_right=0.035, margin_top=0.015, margin_bottom=0.015)
    add_bullet_text(
        s,
        "Intent Object vs Analysis Plan: the Intent Object is what the model said. "
        "The Analysis Plan is what AEGIS was able to bind it to — it adds the pattern, "
        "the join path, the resolved SQL expressions, and the evidence for each binding. "
        "Only the Plan reaches the compiler.",
        Inches(0.95), Inches(5.35), Inches(11.2), Inches(0.8), font_size=12.5)
    add_plain_line(s, "The model said 'average order value'. The definition behind that name was written by the administrator.")

    # 14
    s = add_content_slide("Template-Based SQL Compilation")
    add_stage_tag(s, "Stage 3 of 5")
    add_flat_box(s, "Compiler template", Inches(0.72), Inches(1.42), Inches(3.45), Inches(0.58), light_blue, primary_color, font_size=13)
    add_sql_box(
        s,
        'SELECT {dimension_expr} AS label,\n       {metric_expr} AS value\nFROM {base_table}\n{join_clauses}\nWHERE {mandatory_predicates}\nGROUP BY {dimension_group_expr}\nORDER BY value DESC;',
        Inches(0.78), Inches(2.18), Inches(4.95), Inches(2.58),
        font_size=9.3,
    )
    add_flat_box(s, "Filled from semantic layer", Inches(6.12), Inches(1.42), Inches(2.85), Inches(0.58), soft_green, primary_color, font_size=13)
    add_bullet_text(
        s,
        "- dimension_expr: CASE order status mapping\n- metric_expr: AVG(order total)\n- base_table: Order o\n- joins: none required here\n- predicates: o.Deleted = 0",
        Inches(6.15), Inches(2.12), Inches(2.95), Inches(2.15),
        font_size=11.2,
        line_spacing=0.92,
        min_font_size=9.5,
    )
    add_flat_box(s, "Compiled safe SELECT", Inches(9.42), Inches(1.42), Inches(2.8), Inches(0.58), light_blue, primary_color, font_size=13)
    add_sql_box(
        s,
        'SELECT CASE o.OrderStatusId\n       WHEN 10 THEN "Pending"\n       WHEN 20 THEN "Processing"\n       WHEN 30 THEN "Complete"\n       WHEN 40 THEN "Cancelled"\n       END AS label,\n       AVG(COALESCE(o.OrderTotal, 0)) AS value\nFROM `Order` o\nWHERE o.Deleted = 0\nGROUP BY CASE o.OrderStatusId ... END\nORDER BY value DESC;',
        Inches(9.18), Inches(2.12), Inches(3.15), Inches(2.9),
        font_size=7.2,
    )
    add_bullet_text(s, "The template controls SQL structure; the semantic layer supplies approved identifiers and expressions; user text is never inserted into the query.", Inches(0.9), Inches(5.4), Inches(11.35), Inches(0.62), font_size=13)
    add_plain_line(s, "The query's shape is fixed in advance. Only approved names are dropped into the blanks.")

    # 15 — a full walkthrough on a question that needs four joins, plus the two
    # things a technical reviewer looks for and the earlier slides never showed:
    # where a literal value goes, and what a refusal looks like. Every string on
    # this slide is copied from a real run against the seeded database.
    s = add_content_slide("Walkthrough: A Question That Needs Four Joins")
    add_stage_tag(s, "Stages 1 to 3")
    add_flat_box(s, 'NQ257  "Show product revenue by category."', Inches(0.72), Inches(1.36), Inches(11.6), Inches(0.46), light_blue, primary_color, font_size=13)
    add_flat_box(s, "Intent Object from the model", Inches(0.72), Inches(1.95), Inches(4.3), Inches(0.44), soft_green, primary_color, font_size=12)
    add_sql_box(
        s,
        '{\n  "intent_class": "segment",\n  "metric_term": "line_item_revenue",\n  "dimension_term": "category_name"\n}',
        Inches(0.78), Inches(2.46), Inches(4.2), Inches(1.25), font_size=9.3)
    add_flat_box(s, "Join path resolved by BFS, not by the model", Inches(0.72), Inches(3.85), Inches(4.3), Inches(0.44), soft_green, primary_color, font_size=11.5)
    add_sql_box(
        s,
        'Product -> OrderItem -> Order\n   -> Product_Category_Mapping\n   -> Category',
        Inches(0.78), Inches(4.36), Inches(4.2), Inches(0.92), font_size=9.3)
    add_flat_box(s, "Compiled SQL, executed read-only", Inches(5.35), Inches(1.95), Inches(6.97), Inches(0.44), light_blue, primary_color, font_size=12)
    add_sql_box(
        s,
        'SELECT c.Name AS label, SUM(oi.PriceExclTax) AS value\n'
        'FROM `Order` o\n'
        'LEFT JOIN `OrderItem` oi ON o.Id = oi.OrderId\n'
        'LEFT JOIN `Product` p   ON oi.ProductId = p.Id\n'
        'LEFT JOIN `Product_Category_Mapping` pcm ON p.Id = pcm.ProductId\n'
        'LEFT JOIN `Category` c  ON pcm.CategoryId = c.Id\n'
        'WHERE 1=1\n'
        '  AND o.Deleted = 0\n'
        '  AND p.Deleted = 0\n'
        'GROUP BY c.Name\n'
        'ORDER BY value DESC\n'
        'LIMIT 100',
        Inches(5.4), Inches(2.46), Inches(6.9), Inches(2.82), font_size=8.6)
    tbl = s.shapes.add_table(3, 3, Inches(0.75), Inches(5.42), Inches(11.6), Inches(1.0)).table
    for i, w in enumerate([2.5, 5.05, 4.05]):
        tbl.columns[i].width = Inches(w)
    add_header_row(tbl, ["Also worth seeing", "What actually happens", "Why it matters"], size=11)
    rows = [
        ['"stock less than 10"', "compiles to  p.StockQuantity < @p0   with params {p0: 10}",
         "The number is bound, never pasted into the SQL text"],
        ['"revenue from Google Ads"', "declined: traffic_source is not in the semantic layer",
         "No guess, no nearest-match substitution"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            set_cell(tbl.cell(r, c), value, size=10, bold=(c == 0),
                     fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl, margin_left=0.05, margin_right=0.05, margin_top=0.02, margin_bottom=0.02)

    # 16
    s = add_content_slide("Output Generation")
    fig = os.path.join(figures_dir, 'figure-08-widget-lifecycle.png')
    if os.path.exists(fig):
        add_fit_picture(s, fig, Inches(0.85), Inches(1.45), Inches(5.7), Inches(4.35))
    add_bullet_text(s, "Output is based on result shape, not model narration:\n- KPI: one numeric value with label.\n- Table: row/column data with responsive scrolling in the prototype.\n- Matrix: period columns such as today, week, month, year, and all time.\n- Chart-ready output: categories and values are produced from query results.\n\nThis lets the prototype show both the answer and the evidence path used to produce it.", Inches(6.8), Inches(1.78), Inches(5.45), Inches(4.45), font_size=14.2)

    # 15
    s = add_content_slide("Prototype Implementation")
    add_bullet_text(s, "Prototype stack:\n- FastAPI backend for NL analytics requests.\n- MySQL seeded with nopCommerce-style [7] commerce data.\n- OpenAI-compatible LLM API for intent extraction.\n- Semantic-layer files define metrics, dimensions, filters, and output shapes.\n- Frontend shows answer output plus AEGIS stage evidence.", Inches(0.95), Inches(1.55), Inches(5.65), Inches(3.35), font_size=15.5)
    add_bar_chart(
        s, "Seeded Data Used in Live Prototype",
        ["Order items", "Orders", "Customers", "Products"],
        [6320, 2500, 1200, 17],
        Inches(6.95), Inches(1.55), Inches(5.35), Inches(3.35),
        axis_max=7000, orientation="horizontal")
    add_bullet_text(s, "Reproducibility:\n- Docker build includes app, database seed scripts, and smoke checks.\n- Local health test verified database connection and seeded counts.", Inches(0.95), Inches(5.18), Inches(11.2), Inches(1.0), font_size=14)

    # 17-18 — every figure below is read out of the committed result artifact
    # at build time. Nothing on these two slides is typed by hand, so the deck
    # cannot drift from the evidence: re-run a benchmark, rebuild, and the
    # slide changes with it. A missing artifact renders as "not measured"
    # rather than silently leaving the previous run's number on screen.
    def _metric(payload, key, default="not measured"):
        if not payload:
            return default
        m = (payload.get("metrics") or {}).get(key)
        if not m:
            return default
        return f"{m['n']}/{m['of']} ({m['value']:.1f}%)"

    det = load_result_json('nopcommerce_500_dataset_results.json')
    live = load_result_json('nopcommerce_500_live_benchmark_results.json')
    suite = load_result_json('report_suite_results.json')
    diff = load_result_json('report_differential_results.json')

    s = add_content_slide("Evaluation Track 1: 500 Natural-Language Questions")
    add_bullet_text(
        s,
        "425 questions the semantic layer should answer, 75 realistic e-commerce questions it should decline. "
        "The same corpus is executed two ways, and they answer different questions.",
        Inches(0.85), Inches(1.4), Inches(11.5), Inches(0.6), font_size=13.5)
    tbl = s.shapes.add_table(6, 3, Inches(0.85), Inches(2.1), Inches(11.5), Inches(2.55)).table
    for i, w in enumerate([4.35, 3.4, 3.75]):
        tbl.columns[i].width = Inches(w)
    add_header_row(tbl, ["Measure", "Deterministic run\n(no model in the loop)", "End-to-end run\n(live parser)"], size=11)
    rows = [
        ["Parser produced an intent", "not applicable", _metric(live, "parser_success")],
        ["Supported questions answered", _metric(det, "supported_resolution_validity"), _metric(live, "supported_answer_rate")],
        ["Compiled SQL executed", _metric(det, "supported_execution_validity"), _metric(live, "supported_execution_validity")],
        ["Intent matched the annotation", "not applicable", _metric(live, "supported_intent_exact")],
        ["Out-of-scope questions declined", _metric(det, "boundary_label_validity"), _metric(live, "boundary_rejection_accuracy")],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            set_cell(tbl.cell(r, c), text, size=11, bold=(c == 0),
                     align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER,
                     fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl, margin_left=0.06, margin_right=0.06, margin_top=0.025, margin_bottom=0.025)
    add_bullet_text(
        s,
        "The left column has no model in it: committed intent annotations are fed straight to the mapper, so it "
        "tests the compiler, not the pipeline. Only the right column is an end-to-end result, and it is the one to quote.\n"
        "Declining is only meaningful next to answering: a system that refuses everything scores 100% on the last row.",
        Inches(0.85), Inches(4.8), Inches(11.5), Inches(1.15), font_size=12.5)
    add_plain_line(s, "It answers the questions it was built for, and says no to the ones it was not — on purpose.")

    # Runtime analysis, split at the model boundary. No reference paper in the
    # review reports this split; it is only interesting because the
    # architecture puts SQL authority entirely on one side of it.
    s = add_content_slide("Where the Time Goes")
    def _lat(payload, stage, field="median_ms"):
        st = ((payload or {}).get("latency") or {}).get(stage) or {}
        if not st.get("n"):
            return "not measured"
        return f"{st[field]:,.2f} ms"
    add_bullet_text(
        s,
        "Median and 95th-percentile time per question, measured on the supported questions of the 500-question corpus.",
        Inches(0.85), Inches(1.4), Inches(11.5), Inches(0.45), font_size=13.5)
    tbl = s.shapes.add_table(6, 4, Inches(0.85), Inches(1.98), Inches(11.5), Inches(2.6)).table
    for i, w in enumerate([4.0, 2.5, 2.5, 2.5]):
        tbl.columns[i].width = Inches(w)
    add_header_row(tbl, ["Stage", "Who runs it", "Median", "95th percentile"], size=11.5)
    rows = [
        ["1. Intent extraction", "The language model", _lat(live, "parse_ms"), _lat(live, "parse_ms", "p95_ms")],
        ["2. Semantic resolution", "AEGIS", _lat(det, "resolve_ms"), _lat(det, "resolve_ms", "p95_ms")],
        ["3. SQL compilation", "AEGIS", _lat(det, "compile_ms"), _lat(det, "compile_ms", "p95_ms")],
        ["4. Database execution", "MySQL", _lat(det, "execute_ms"), _lat(det, "execute_ms", "p95_ms")],
        ["Stages 2-4 combined", "Everything after the model", _lat(live, "deterministic_ms"), _lat(live, "deterministic_ms", "p95_ms")],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            set_cell(tbl.cell(r, c), text, size=11.5, bold=(c == 0 or r == 5),
                     align=PP_ALIGN.LEFT if c < 2 else PP_ALIGN.CENTER,
                     fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl, margin_left=0.07, margin_right=0.07, margin_top=0.03, margin_bottom=0.03)
    add_bullet_text(
        s,
        "Almost all of the wall clock is the model reading the question. The part that decides which tables to "
        "join, writes the SQL, and runs it takes a few milliseconds — and the model has no influence over any of it.\n"
        "The model figure is a property of the gateway used here, not of the architecture; the deterministic "
        "figures are the ones that would follow AEGIS to another deployment.",
        Inches(0.85), Inches(4.75), Inches(11.5), Inches(1.35), font_size=12.5)

    s = add_content_slide("Evaluation Track 2: nopCommerce's Own 20 Admin Reports")
    add_bullet_text(
        s,
        "The report list is nopCommerce's own admin menu, and the comparison target is nopCommerce's own "
        "report code, read from source at commit 64bdf2ff. Neither was chosen by this project.",
        Inches(0.85), Inches(1.38), Inches(11.5), Inches(0.6), font_size=13.5)
    suite_line = (f"{suite['reproduced']}/{suite['total']}" if suite else "not measured")
    diff_line = (f"{diff['matched']}/{diff['total']} ({diff['accuracy']:.1f}%)" if diff else "not measured")
    tbl = s.shapes.add_table(3, 3, Inches(0.85), Inches(2.08), Inches(11.5), Inches(1.35)).table
    for i, w in enumerate([3.6, 2.2, 5.7]):
        tbl.columns[i].width = Inches(w)
    add_header_row(tbl, ["Check", "Result", "What it does and does not show"], size=11.5)
    rows = [
        ["Report reached an answer\nand compiled to SQL", suite_line,
         "Coverage of the report shape. Not correctness — several once passed this while silently wrong."],
        ["Result set matched the\nplatform's own query", diff_line,
         "The check that tests the claim: same database, same rows, same values."],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            set_cell(tbl.cell(r, c), text, size=11, bold=(c == 1),
                     align=PP_ALIGN.CENTER if c == 1 else PP_ALIGN.LEFT,
                     fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl, margin_left=0.06, margin_right=0.06, margin_top=0.025, margin_bottom=0.025)
    add_flat_box(s, "The five reports that did not match, and why", Inches(0.85), Inches(3.62), Inches(11.5), Inches(0.44), soft_green, primary_color, font_size=12.5)
    tbl = s.shapes.add_table(4, 2, Inches(0.85), Inches(4.16), Inches(11.5), Inches(1.5)).table
    tbl.columns[0].width = Inches(4.6)
    tbl.columns[1].width = Inches(6.9)
    add_header_row(tbl, ["Difference", "Reports affected"], size=11)
    rows = [
        ["Row count only — the platform's report\ncarries its own LIMIT (5, 15, 100)",
         "Bestsellers by quantity, Bestsellers by amount, Latest orders, Best customers by order count"],
        ["Label column — customer name where the\nplatform's query labels by email",
         "Best customers by order total, Best customers by order count"],
        ["Wrong value", "None. Every overlapping row agrees to the last decimal."],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            set_cell(tbl.cell(r, c), text, size=10.5, bold=(r == 3),
                     fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl, margin_left=0.06, margin_right=0.06, margin_top=0.02, margin_bottom=0.02)
    add_plain_line(s, "Where the numbers differ from nopCommerce's own reports, they do not — only how many rows and which label.")

    # 18
    s = add_content_slide("Beneficiaries and Expected Impact")
    tbl = s.shapes.add_table(5, 3, Inches(0.9), Inches(1.55), Inches(11.55), Inches(4.55)).table
    widths = [2.9, 4.0, 4.65]
    for i, w in enumerate(widths):
        tbl.columns[i].width = Inches(w)
    add_header_row(tbl, ["Beneficiary", "Current difficulty", "AEGIS impact"], size=12)
    rows = [
        ["Business users", "Need reports but need new report combinations", "Ask natural-language analytics questions and receive controlled table or chart-ready results"],
        ["Developers / BI teams", "Repeated report requests consume implementation time", "Reusable semantic definitions reduce one-off SQL writing for common analytics"],
        ["Database administrators", "Model-written SQL is hard to audit and restrict", "Only approved metrics, dimensions, joins, and read-only SQL reach the database"],
        ["Researchers", "NL2SQL work often mixes language accuracy with execution authority", "Shows an architecture where language understanding and SQL authority are separated"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            set_cell(tbl.cell(r, c), text, size=11.5, bold=(c == 0), fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl, margin_left=0.06, margin_right=0.06, margin_top=0.025, margin_bottom=0.025)
    add_bullet_text(s, "Expected impact: safer self-service analytics for repeated e-commerce reporting, without giving the LLM permission to author executable SQL.", Inches(1.05), Inches(6.18), Inches(11.1), Inches(0.45), font_size=12.5)

    # 19
    s = add_content_slide("Scope and Limitations")
    add_bullet_text(
        s,
        "Current evaluation scope:\n"
        "- Implemented and evaluated for nopCommerce [7] analytics on MySQL; the semantic layer is deliberately finite and deployment-specific.\n"
        "- Because the layer is per-deployment, cross-domain benchmarks such as Spider and BIRD do not apply: they test generalisation to unseen schemas, which this architecture does not attempt.\n\n"
        "Limitations I would raise before the committee does:\n"
        "- Correctness at scale is measured on the 20 reports, not on all 425 supported questions. For those, I show that the SQL runs — not that every answer is right.\n"
        "- The gateway used for intent extraction resolves a model alias per request, so the run is not pinned to a single model. Each result row records what actually served it.\n"
        "- Five reports differ from the platform's own output in row count or label column. No report differs in value.\n"
        "- Vague, unsupported language can still be misread during intent extraction.\n"
        "- Other SQL dialects require compiler extensions.\n\n"
        "Not a thesis limitation:\n"
        "- Multi-turn conversation is out of scope because the work evaluates single-request analytics.",
        Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.2), font_size=13.5)

    # 20
    s = add_content_slide("Conclusion")
    add_bullet_text(s, "Conclusion:\n- AEGIS keeps the LLM useful for language understanding while removing SQL-authoring authority from the model.\n- The semantic layer makes analytics definitions explicit, auditable, and reusable for a target system.\n- The compiler produces read-only SQL from approved definitions and blocks unsupported execution paths.\n- The nopCommerce [7] prototype, the 500-question corpus, and the differential against the platform's own report logic show that the approach can be implemented and evaluated practically.\n\nMain contribution:\n- A constraint-based architecture for safe LLM-assisted natural language analytics.", Inches(1.2), Inches(1.65), Inches(10.9), Inches(4.85), font_size=17)

    # 21
    s = add_content_slide("Future Work")
    add_bullet_text(s, "Possible extensions:\n- Evaluate AEGIS on another e-commerce system by replacing only the semantic layer.\n- Add more oracle reports from real admin and business workflows.\n- Extend compiler modules for PostgreSQL and SQL Server.\n- Improve clarification behavior for ambiguous but answerable user questions.\n- Add richer chart recommendation while keeping SQL generation deterministic.", Inches(1.25), Inches(1.8), Inches(10.7), Inches(4.35), font_size=18)

    # 22
    s = add_content_slide("References")
    refs_box = s.shapes.add_textbox(Inches(0.8), Inches(1.78), Inches(11.75), Inches(4.8))
    refs_tf = refs_box.text_frame
    refs_tf.word_wrap = True
    references = [
        "[1] Li, F., & Jagadish, H. V. (2014). Constructing an interactive natural language interface for relational databases. PVLDB.",
        "[2] Narechania, A., Srinivasan, A., & Stasko, J. (2021). nl4dv: A toolkit for generating analytic specifications from natural language. IEEE TVCG.",
        "[3] Deng, D., Wu, A., Qu, H., & Wu, Y. (2023). DashBot: Insight-driven dashboard generation. IEEE TVCG.",
        "[4] Scholak, T., Schucher, N., & Bahdanau, D. (2021). PICARD: Parsing incrementally for constrained auto-regressive decoding. EMNLP.",
        "[5] Shalaan, H. S. et al. (2025). G-SQL: A schema-aware and rule-guided approach for robust natural language to SQL translation. IEEE Access.",
        "[6] Su, X. et al. (2026). A robust natural language text-to-SQL generation framework with dynamic strategies based on LLMs. Scientific Reports.",
        "[7] nopSolutions. nopCommerce open-source e-commerce platform. GitHub repository.",
    ]
    for i, ref_text in enumerate(references):
        p = refs_tf.paragraphs[0] if i == 0 else refs_tf.add_paragraph()
        run = p.add_run()
        run.text = ref_text
        run.font.name = TEMPLATE_FONT
        run.font.size = Pt(12)
        run.font.color.rgb = BODY_COLOR
        _set_hanging_indent(p, 0.3, -0.3)
        p.space_after = Pt(9)
        p.line_spacing = 1.05

    # 23
    s = add_content_slide("Thank You & Discussion")
    qa = s.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(2))
    qtf = qa.text_frame
    p1 = qtf.paragraphs[0]
    p1.text = "THANK YOU!"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(36)
    p1.font.name = TEMPLATE_FONT
    p1.font.bold = True
    p1.font.color.rgb = primary_color
    p2 = qtf.add_paragraph()
    p2.text = "Questions & Final Defense Discussion"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(20)
    p2.font.name = TEMPLATE_FONT
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(70, 70, 70)

    # Appendix — shown only if asked. Every block below is the project's own
    # source, pasted verbatim rather than paraphrased, because the fastest way
    # to answer "did you actually build this?" is to show the file.
    s = add_content_slide("Appendix A: The Semantic Layer Is a File I Wrote")
    add_bullet_text(s, "aegis/server/semantic_layer.py — two of the objects the model is allowed to name:",
                    Inches(0.85), Inches(1.36), Inches(11.5), Inches(0.42), font_size=13)
    add_sql_box(
        s,
        'Metric(\n'
        '    id="avg_order_value",\n'
        '    label="Average Order Value",\n'
        '    description="Average revenue per order",\n'
        '    sql_expr="AVG(COALESCE(o.OrderTotal, 0))",\n'
        '    binding_table="Order",\n'
        '    required_joins=[],\n'
        '    default_visual="kpi_card",\n'
        '    time_anchor="o.CreatedOnUtc",\n'
        ')',
        Inches(0.85), Inches(1.9), Inches(5.6), Inches(2.5), font_size=9.5)
    add_sql_box(
        s,
        'Dimension(\n'
        '    id="order_status",\n'
        '    label="Order Status",\n'
        '    sql_expr="CASE o.OrderStatusId "\n'
        '             "WHEN 10 THEN \'Pending\' "\n'
        '             "WHEN 20 THEN \'Processing\' "\n'
        '             "WHEN 30 THEN \'Complete\' "\n'
        '             "WHEN 40 THEN \'Cancelled\' "\n'
        '             "ELSE \'Unknown\' END",\n'
        '    binding_table="Order",\n'
        '    datatype="string",\n'
        ')',
        Inches(6.72), Inches(1.9), Inches(5.6), Inches(2.5), font_size=9.5)
    add_bullet_text(
        s,
        "The model may return the string \"avg_order_value\". It cannot return the expression, the table, or the join. "
        "Deploying AEGIS on another system means writing this file again — and nothing else.",
        Inches(0.85), Inches(4.55), Inches(11.5), Inches(0.9), font_size=13)
    add_plain_line(s, "This list is the whole vocabulary. Anything not written here cannot be asked for.")

    s = add_content_slide("Appendix B: The Compiler Fills Blanks, It Does Not Write SQL")
    add_bullet_text(s, "aegis/server/compiler.py — the function that builds the SELECT clause:",
                    Inches(0.85), Inches(1.36), Inches(11.5), Inches(0.42), font_size=13)
    add_sql_box(
        s,
        'def _assemble_select(self, metric, dimension, extra_metrics=None, plan=None) -> str:\n'
        '    measures = [f"{metric.sql_expr} AS value"]\n'
        '    for extra in extra_metrics or []:\n'
        '        measures.append(f"{extra.sql_expr} AS `{extra.id}`")\n'
        '    projected = ", ".join(measures)\n'
        '\n'
        '    if dimension:\n'
        '        dim_expr = self._dimension_select_expr(dimension, plan)\n'
        '        return f"SELECT {dim_expr} AS label, {projected}"\n'
        '    return f"SELECT {projected}"\n'
        '\n'
        '# Every interpolated value is an attribute of a semantic-layer object.\n'
        '# No argument to this function carries user text.',
        Inches(0.85), Inches(1.9), Inches(11.5), Inches(3.0), font_size=10)
    add_bullet_text(
        s,
        "Literal values take a separate route: a filter compiles to a placeholder and a bound parameter "
        "(p.StockQuantity < @p0 with {p0: 10}), so a value never becomes part of the SQL text. "
        "After assembly, a forbidden-pattern scan rejects any non-SELECT construct before execution.",
        Inches(0.85), Inches(5.05), Inches(11.5), Inches(0.95), font_size=12.5)

    s = add_content_slide("Appendix C: Why Templates Instead of Constrained Decoding")
    tbl = s.shapes.add_table(4, 3, Inches(0.85), Inches(1.5), Inches(11.5), Inches(3.1)).table
    for i, w in enumerate([2.6, 4.45, 4.45]):
        tbl.columns[i].width = Inches(w)
    add_header_row(tbl, ["", "PICARD [4] — constrained decoding", "AEGIS — template compilation"], size=12)
    rows = [
        ["What the model emits", "SQL text, token by token", "A typed Intent Object with no SQL field"],
        ["What is checked", "Each token against a parser, during generation", "Each id against the semantic layer, before compilation"],
        ["What we give up", "Nothing — any valid SQL stays reachable", "Flexibility: a question outside the semantic layer cannot be answered at all"],
    ]
    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            set_cell(tbl.cell(r, c), text, size=12, bold=(c == 0),
                     fill=(RGBColor(0xF5, 0xF7, 0xFB) if r % 2 == 0 else None))
    style_table(tbl, margin_left=0.07, margin_right=0.07, margin_top=0.03, margin_bottom=0.03)
    add_bullet_text(
        s,
        "Constrained decoding guarantees the SQL parses. It does not guarantee the SQL is the query the "
        "administrator would have approved — a syntactically perfect query can still join the wrong grain, "
        "drop a soft-delete filter, or read a column nobody meant to expose.\n\n"
        "The trade is real and worth stating plainly: AEGIS buys auditability by giving up open-ended coverage. "
        "That is why 75 of the 500 evaluation questions are ones it must decline.",
        Inches(0.85), Inches(4.75), Inches(11.5), Inches(1.35), font_size=13)

    footer_names = ('Date Placeholder 3', 'Footer Placeholder 4', 'Slide Number Placeholder 5',
                    'Rectangle 2', 'Rectangle 6', 'Rectangle 10', 'TextBox 14')
    band_top_emu = Inches(FOOTER_BAND_TOP_INCHES)
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.name in footer_names or shape.top is None or shape.height is None:
                continue
            bottom = shape.top + shape.height
            if bottom > band_top_emu:
                OVERFLOW_WARNINGS.append(
                    "slide %d: shape %r extends to %.2f\", past the %.2f\" footer band"
                    % (idx, shape.name, bottom / 914400, FOOTER_BAND_TOP_INCHES))

    prs.save(output_path)
    print(f"Successfully generated final-defense presentation: {output_path}")
    if OVERFLOW_WARNINGS:
        print(f"\n{len(OVERFLOW_WARNINGS)} layout warning(s):")
        for warning in OVERFLOW_WARNINGS:
            print(f"  - {warning}")
    else:
        print("Layout check: no content overflows the footer band.")


if __name__ == '__main__':
    create_presentation()
