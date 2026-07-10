"""
Generate AEGIS Thesis Defense Presentation (.pptx)
Run from repo root: python3 docs/scripts/generate_presentation.py
Output: docs/AEGIS_Thesis_Defense.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Color Palette ────────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)   # dark navy  - headings/BG
TEAL       = RGBColor(0x1A, 0x73, 0xA0)   # steel blue - accents
GOLD       = RGBColor(0xF1, 0x8F, 0x01)   # amber      - highlights
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHTBG    = RGBColor(0xF0, 0xF4, 0xF8)   # slide background
DARK       = RGBColor(0x1A, 0x1A, 0x2E)   # body text
MID        = RGBColor(0x44, 0x5B, 0x7A)   # subtext
GREENOK    = RGBColor(0x27, 0xAE, 0x60)
REDNG      = RGBColor(0xE7, 0x4C, 0x3C)
LIGHTBLUE  = RGBColor(0xD6, 0xEA, 0xF8)
LIGHTYELLOW= RGBColor(0xFE, 0xF9, 0xE7)
LIGHTGRAY  = RGBColor(0xEC, 0xF0, 0xF1)
TEALLIGHT  = RGBColor(0xD0, 0xE8, 0xF3)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helper utilities ─────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=None, border=None):
    # 'border' is an alias for 'line' for readability at call sites
    line = line or border
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def txb(slide, text, l, t, w, h,
        size=18, bold=False, italic=False, color=DARK,
        align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def txb_lines(slide, lines, l, t, w, h,
              size=16, bold=False, color=DARK,
              align=PP_ALIGN.LEFT, spacing=1.15, font="Calibri"):
    """Multiple lines in one textbox with line spacing."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (line_text, line_bold, line_size, line_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', f'{int(spacing*100)}%')
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', '0')
        run = p.add_run()
        run.text = line_text
        run.font.name = font
        run.font.size = _Pt(line_size or size)
        run.font.bold = line_bold if line_bold is not None else bold
        run.font.color.rgb = line_color or color
    return tb

def bullet_box(slide, items, l, t, w, h,
               size=16, color=DARK, bullet_color=TEAL,
               indent=0.18, spacing=115, font="Calibri"):
    """Bullet list with colored square bullets."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    from pptx.oxml.ns import qn
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, sz, clr, is_sub = item
        else:
            text, sz, clr, is_sub = item, size, color, False

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT

        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', f'{spacing}%')
        spcBef = etree.SubElement(pPr, qn('a:spcBef'))
        spcPts = etree.SubElement(spcBef, qn('a:spcPts'))
        spcPts.set('val', '60' if is_sub else '100')

        bullet_indent = indent * 1.8 if is_sub else indent
        run_bullet = p.add_run()
        run_bullet.text = "  • " if is_sub else "▸  "
        run_bullet.font.name = "Segoe UI Symbol"
        run_bullet.font.size = Pt(sz - 2 if is_sub else sz)
        run_bullet.font.color.rgb = MID if is_sub else bullet_color
        run_bullet.font.bold = False

        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(sz)
        run.font.color.rgb = clr
        run.font.bold = False
    return tb

def header_bar(slide, title, subtitle=None):
    """Top navy header bar used on most content slides."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), fill=NAVY)
    add_rect(slide, 0, Inches(1.25), SLIDE_W, Inches(0.06), fill=GOLD)
    txb(slide, title,
        Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.7),
        size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri Light")
    if subtitle:
        txb(slide, subtitle,
            Inches(0.45), Inches(0.82), Inches(12.4), Inches(0.4),
            size=16, color=TEALLIGHT, align=PP_ALIGN.LEFT)

def slide_bg(slide, color=WHITE):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=color)

def footer(slide, text="Md. Riaz  |  Pundra University of Science and Technology  |  BSc CSE Thesis Defense"):
    add_rect(slide, 0, Inches(7.15), SLIDE_W, Inches(0.35), fill=NAVY)
    txb(slide, text,
        Inches(0.3), Inches(7.18), Inches(12.7), Inches(0.28),
        size=9.5, color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)

def card(slide, l, t, w, h, fill=LIGHTBG, border=None):
    r = add_rect(slide, l, t, w, h, fill=fill, line=border or TEAL,
                 line_w=Pt(1))
    return r

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, NAVY)
# decorative side stripe
add_rect(s, 0, 0, Inches(0.55), SLIDE_H, fill=GOLD)
add_rect(s, Inches(0.55), 0, Inches(0.08), SLIDE_H, fill=TEAL)

# university / degree badge
add_rect(s, Inches(1.1), Inches(0.35), Inches(11.7), Inches(0.55),
         fill=RGBColor(0x0D, 0x1B, 0x33))
txb(s, "PUNDRA UNIVERSITY OF SCIENCE AND TECHNOLOGY  ·  DEPARTMENT OF CSE  ·  8TH SEMESTER THESIS DEFENSE",
    Inches(1.2), Inches(0.38), Inches(11.5), Inches(0.5),
    size=10.5, color=GOLD, align=PP_ALIGN.CENTER, font="Calibri")

# AEGIS acronym
txb(s, "AEGIS",
    Inches(1.1), Inches(1.1), Inches(11.6), Inches(1.4),
    size=96, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri Light")

# gold divider
add_rect(s, Inches(2.5), Inches(2.55), Inches(8.3), Inches(0.055), fill=GOLD)

# full title
txb(s, "A Safety-by-Design Architecture for\nLLM-Driven Self-Service Analytics",
    Inches(1.1), Inches(2.65), Inches(11.6), Inches(1.1),
    size=28, bold=False, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri Light")

# acronym expansion
txb(s, "Analytics Engine with Guaranteed Injection Safety",
    Inches(1.1), Inches(3.78), Inches(11.6), Inches(0.45),
    size=16, italic=True, color=TEALLIGHT, align=PP_ALIGN.CENTER)

# gold divider 2
add_rect(s, Inches(4.0), Inches(4.28), Inches(5.3), Inches(0.035), fill=GOLD)

# author + institution
txb(s, "Md. Riaz",
    Inches(1.1), Inches(4.4), Inches(11.6), Inches(0.45),
    size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "Pundra University of Science and Technology",
    Inches(1.1), Inches(4.88), Inches(11.6), Inches(0.4),
    size=15, color=TEALLIGHT, align=PP_ALIGN.CENTER)
txb(s, "Bachelor of Science in Computer Science & Engineering  |  2026",
    Inches(1.1), Inches(5.3), Inches(11.6), Inches(0.35),
    size=13.5, color=RGBColor(0x8A, 0xA8, 0xC8), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "Presentation Outline")
footer(s)

items_left = [
    ("The Problem — Why This Research?", False, 20, NAVY, False),
    ("Research Objectives & Questions", False, 20, NAVY, False),
    ("What Makes AEGIS Novel?", False, 20, NAVY, False),
    ("System Architecture — 7-Stage Pipeline", False, 20, NAVY, False),
    ("The Semantic Layer — Core Contribution", False, 20, NAVY, False),
]
items_right = [
    ("How a Query Flows Through AEGIS", False, 20, NAVY, False),
    ("SQL Safety — Formal Guarantee", False, 20, NAVY, False),
    ("Evaluation Results (100-Query Benchmark)", False, 20, NAVY, False),
    ("Ablation Study & Generalizability", False, 20, NAVY, False),
    ("Conclusion & Future Work", False, 20, NAVY, False),
]

for i, (txt, bold, sz, clr, _) in enumerate(items_left):
    y = Inches(1.55) + i * Inches(0.88)
    add_rect(s, Inches(0.5), y, Inches(0.48), Inches(0.48), fill=TEAL)
    txb(s, str(i+1), Inches(0.5), y, Inches(0.48), Inches(0.48),
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, txt, Inches(1.1), y + Inches(0.04), Inches(5.4), Inches(0.48),
        size=17, color=DARK, align=PP_ALIGN.LEFT)

for i, (txt, bold, sz, clr, _) in enumerate(items_right):
    y = Inches(1.55) + i * Inches(0.88)
    add_rect(s, Inches(7.0), y, Inches(0.48), Inches(0.48), fill=GOLD)
    txb(s, str(i+6), Inches(7.0), y, Inches(0.48), Inches(0.48),
        size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, txt, Inches(7.6), y + Inches(0.04), Inches(5.4), Inches(0.48),
        size=17, color=DARK, align=PP_ALIGN.LEFT)

add_rect(s, Inches(6.55), Inches(1.45), Inches(0.05), Inches(5.3), fill=LIGHTGRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "The Problem", "Why do we need a new approach to analytics?")
footer(s)

# left pain point cards
pains = [
    (REDNG,  "SQL Barrier",      "Business users cannot write SQL.\nThey wait days for a developer\nto build each report."),
    (REDNG,  "Injection Risk",   "LLM-generated SQL can be\nmanipulated into DROP TABLE or\ndata exfiltration by a bad prompt."),
    (REDNG,  "Hallucination",    "LLMs invent column names and\njoin conditions that produce\nsilently wrong answers."),
    (REDNG,  "One-Shot Results", "Existing tools return a one-time\nquery result — not a refreshable,\nreusable dashboard widget."),
]
for i, (clr, title, desc) in enumerate(pains):
    row, col = divmod(i, 2)
    x = Inches(0.4) + col * Inches(6.2)
    y = Inches(1.5) + row * Inches(2.45)
    card(s, x, y, Inches(5.8), Inches(2.25), fill=RGBColor(0xFF,0xF0,0xF0), border=clr)
    add_rect(s, x, y, Inches(0.18), Inches(2.25), fill=clr)
    txb(s, title, x + Inches(0.28), y + Inches(0.12),
        Inches(5.3), Inches(0.5), size=18, bold=True, color=clr)
    txb(s, desc,  x + Inches(0.28), y + Inches(0.6),
        Inches(5.3), Inches(1.5), size=14.5, color=DARK)

# stat callout bottom
add_rect(s, Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.5), fill=NAVY)
txb(s, "61% of business reporting questions are recurring — the same report with a different date range or department.",
    Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
    size=13, color=GOLD, align=PP_ALIGN.CENTER, bold=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — RESEARCH OBJECTIVES
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "Research Objectives & Questions")
footer(s)

rqs = [
    ("RQ1", "Accuracy",        "Can a structured semantic layer guide the LLM to correctly interpret natural-language analytics queries with ≥95% accuracy?"),
    ("RQ2", "Safety",          "Does structural prevention (template-based compilation) eliminate SQL injection more reliably than detection-based approaches?"),
    ("RQ3", "Ablation",        "How much does each pipeline component (vocabulary injection, BFS joins, safety scanner) contribute to overall accuracy?"),
    ("RQ4", "Generalizability","Can AEGIS be adapted to a different e-commerce schema (WooCommerce) without pipeline changes?"),
    ("RQ5", "Latency",         "Does the safety-first architecture introduce unacceptable overhead for interactive use?"),
]

for i, (tag, label, question) in enumerate(rqs):
    y = Inches(1.45) + i * Inches(1.06)
    # tag box
    add_rect(s, Inches(0.35), y, Inches(0.85), Inches(0.78), fill=NAVY)
    txb(s, tag, Inches(0.35), y + Inches(0.03), Inches(0.85), Inches(0.42),
        size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txb(s, label, Inches(0.35), y + Inches(0.42), Inches(0.85), Inches(0.33),
        size=9.5, color=TEALLIGHT, align=PP_ALIGN.CENTER)
    # question text
    add_rect(s, Inches(1.3), y + Inches(0.1), Inches(11.5), Inches(0.6),
             fill=LIGHTBG, border=LIGHTGRAY)
    txb(s, question, Inches(1.45), y + Inches(0.15), Inches(11.2), Inches(0.55),
        size=15, color=DARK)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — NOVELTY / CORE THESIS
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "Core Novelty — Structural Safety vs. Detection", "What makes AEGIS fundamentally different from every prior NL2SQL system?")
footer(s)

# Left: Prior approach
card(s, Inches(0.35), Inches(1.45), Inches(5.85), Inches(5.3),
     fill=RGBColor(0xFF,0xF5,0xF5), border=REDNG)
add_rect(s, Inches(0.35), Inches(1.45), Inches(5.85), Inches(0.52), fill=REDNG)
txb(s, "Prior Approaches (End-to-End NL2SQL)",
    Inches(0.45), Inches(1.5), Inches(5.65), Inches(0.42),
    size=14.5, bold=True, color=WHITE)
txb(s, "User Query  →  LLM  →  SQL String  →  Execute",
    Inches(0.5), Inches(2.1), Inches(5.6), Inches(0.4),
    size=13.5, italic=True, color=REDNG, bold=True)
prior_pts = [
    "LLM sees the full database schema",
    "LLM writes raw SQL — any output is possible",
    "Injection caught (or missed) by post-hoc filters",
    "Hallucinated columns fail silently at runtime",
    "Safety = detection — you can always miss one",
]
bullet_box(s, [(p, 13.5, DARK, False) for p in prior_pts],
           Inches(0.5), Inches(2.55), Inches(5.55), Inches(3.0))

# Right: AEGIS approach
card(s, Inches(7.1), Inches(1.45), Inches(5.85), Inches(5.3),
     fill=RGBColor(0xF0,0xFF,0xF5), border=GREENOK)
add_rect(s, Inches(7.1), Inches(1.45), Inches(5.85), Inches(0.52), fill=GREENOK)
txb(s, "AEGIS (This Thesis)",
    Inches(7.2), Inches(1.5), Inches(5.65), Inches(0.42),
    size=14.5, bold=True, color=WHITE)
txb(s, "User Query  →  LLM  →  Intent JSON  →  Compiler  →  SQL",
    Inches(7.2), Inches(2.1), Inches(5.6), Inches(0.4),
    size=13.5, italic=True, color=GREENOK, bold=True)
aegis_pts = [
    "LLM never sees the database schema",
    "LLM outputs structured JSON — no SQL output channel",
    "SQL generated only from pre-approved templates",
    "Column names are compile-time constants, not LLM output",
    "Safety = prevention — injection is structurally impossible",
]
bullet_box(s, [(p, 13.5, DARK, False) for p in aegis_pts],
           Inches(7.2), Inches(2.55), Inches(5.55), Inches(3.0))

# VS divider
add_rect(s, Inches(6.3), Inches(2.5), Inches(0.7), Inches(0.7), fill=GOLD)
txb(s, "VS", Inches(6.3), Inches(2.53), Inches(0.7), Inches(0.65),
    size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Bottom key insight
add_rect(s, Inches(0.35), Inches(6.82), Inches(12.6), Inches(0.5), fill=NAVY)
txb(s, "Key Insight: You cannot make LLM-generated SQL safe by checking it. You make it safe by ensuring it is never generated.",
    Inches(0.5), Inches(6.87), Inches(12.3), Inches(0.4),
    size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "System Architecture — 7-Stage Pipeline")
footer(s)

stages = [
    ("1", "Intent\nExtraction",   "LLM maps natural language\nto structured IntentObject",    TEAL,  "intent_parser.py"),
    ("2", "Coverage\nValidation", "Rejects queries outside\nthe semantic vocabulary",          NAVY,  "run_demo_server.py"),
    ("3", "Semantic\nMapping",    "Resolves terms to canonical\nIDs, expands business logic",  TEAL,  "mapper.py"),
    ("4", "Permission\nRewriting","Appends row-level security\nWHERE clause predicates",       NAVY,  "permission_rewriter.py"),
    ("5", "SQL\nCompilation",     "BFS join resolution +\ntemplate-based SQL generation",      TEAL,  "compiler.py"),
    ("6", "Visualization\nSelect","Rule-based chart type\nselection (11 primitives)",          NAVY,  "visualization.py"),
    ("7", "Widget\nPersistence",  "SHA-256 dedup, storage\nand scheduled refresh",             TEAL,  "widget_engine.py"),
]

# LLM zone label
add_rect(s, Inches(0.35), Inches(1.4), Inches(2.35), Inches(5.55),
         fill=RGBColor(0xFF,0xF9,0xF0), border=GOLD)
txb(s, "LLM\nZone", Inches(0.38), Inches(1.5), Inches(2.28), Inches(0.6),
    size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Compiler zone label
add_rect(s, Inches(2.8), Inches(1.4), Inches(10.2), Inches(5.55),
         fill=RGBColor(0xF0,0xF8,0xFF), border=TEAL)
txb(s, "Deterministic Compiler Zone  (LLM output never reaches here as raw text)",
    Inches(2.85), Inches(1.5), Inches(10.1), Inches(0.45),
    size=11, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

for i, (num, title, desc, clr, module) in enumerate(stages):
    x = Inches(0.45) + i * Inches(1.82)
    y = Inches(2.1)
    bh = Inches(4.5)
    bw = Inches(1.65)
    card(s, x, y, bw, bh, fill=WHITE, border=clr)
    add_rect(s, x, y, bw, Inches(0.42), fill=clr)
    txb(s, num, x, y + Inches(0.04), bw, Inches(0.38),
        size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, title, x + Inches(0.06), y + Inches(0.5), bw - Inches(0.12), Inches(0.75),
        size=13, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txb(s, desc,  x + Inches(0.08), y + Inches(1.3), bw - Inches(0.16), Inches(2.0),
        size=11, color=DARK, align=PP_ALIGN.CENTER)
    add_rect(s, x + Inches(0.08), y + Inches(3.5), bw - Inches(0.16), Inches(0.55),
             fill=LIGHTBG, border=clr)
    txb(s, module, x + Inches(0.1), y + Inches(3.55), bw - Inches(0.2), Inches(0.45),
        size=9, color=clr, align=PP_ALIGN.CENTER, italic=True)
    # arrow
    if i < 6:
        ax = x + bw + Inches(0.03)
        ay = y + Inches(1.9)
        add_rect(s, ax, ay, Inches(0.14), Inches(0.14), fill=GOLD)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — SEMANTIC LAYER OVERVIEW
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "The Semantic Layer — Core Contribution",
           "The closed vocabulary that makes structural safety possible")
footer(s)

# Central concept box
add_rect(s, Inches(3.5), Inches(1.45), Inches(6.35), Inches(1.0), fill=NAVY)
txb(s, "The semantic layer defines the COMPLETE, ENUMERABLE set",
    Inches(3.55), Inches(1.52), Inches(6.25), Inches(0.42),
    size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s, "of things AEGIS can answer — everything else is rejected.",
    Inches(3.55), Inches(1.88), Inches(6.25), Inches(0.42),
    size=14, color=TEALLIGHT, align=PP_ALIGN.CENTER)

# 4 component cards
components = [
    (TEAL,  "15 Metrics",         "Named SQL aggregate expressions.\nLLM picks an ID — compiler inserts\nthe pre-approved SQL fragment.",
     "revenue, order_count,\navg_order_value, item_quantity,\nprofit, shipment_count…"),
    (NAVY,  "34 Dimensions",      "Grouping & filtering axes.\nCovers time, geography, product,\ncustomer, order status, shipment.",
     "order_month, category_name,\ncustomer_email, order_status,\nbilling_city, store_name…"),
    (GOLD,  "11 JOIN Paths",      "Pre-approved ON clauses.\nCompiler runs BFS to find the\nminimal join path automatically.",
     "Order→Customer→Address\nOrder→OrderItem→Product\nProduct→Category…"),
    (TEAL,  "Vocabulary Injection","All metric+dimension labels\ninjected into LLM system prompt.\nSYNONYMS = {} intentionally.",
     "~1,100 tokens of approved\nvocabulary. LLM maps 'sales'\nto 'revenue' at inference time."),
]

for i, (clr, title, desc, example) in enumerate(components):
    col, row = divmod(i, 2)
    x = Inches(0.35) + col * Inches(6.5)
    y = Inches(2.65) + row * Inches(2.25)
    card(s, x, y, Inches(6.15), Inches(2.1), fill=WHITE, border=clr)
    add_rect(s, x, y, Inches(6.15), Inches(0.45), fill=clr)
    txb(s, title, x + Inches(0.15), y + Inches(0.06),
        Inches(5.85), Inches(0.38), size=17, bold=True, color=WHITE)
    txb(s, desc,  x + Inches(0.15), y + Inches(0.55),
        Inches(3.1), Inches(1.45), size=12.5, color=DARK)
    add_rect(s, x + Inches(3.3), y + Inches(0.52), Inches(2.7), Inches(1.45),
             fill=LIGHTBG, border=clr)
    txb(s, example, x + Inches(3.4), y + Inches(0.6),
        Inches(2.55), Inches(1.3), size=11, color=MID, italic=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — 15 METRICS DEEP DIVE
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "Semantic Layer — 15 Metrics",
           "Every aggregate SQL expression is a fixed, pre-approved constant — the LLM only picks the ID")
footer(s)

metrics = [
    ("revenue",           "SUM(COALESCE(o.OrderTotal,0) - COALESCE(o.RefundedAmount,0))"),
    ("order_count",       "COUNT(DISTINCT o.Id)"),
    ("avg_order_value",   "AVG(COALESCE(o.OrderTotal, 0))"),
    ("item_quantity",     "SUM(COALESCE(oi.Quantity, 0))"),
    ("shipping_cost",     "SUM(o.OrderShippingExclTax)"),
    ("customer_count",    "COUNT(DISTINCT cu.Id)"),
    ("refund_count",      "COUNT(DISTINCT CASE WHEN o.RefundedAmount > 0 THEN o.Id END)"),
    ("refund_amount",     "SUM(o.RefundedAmount)"),
    ("discount_amount",   "SUM(o.OrderDiscount)"),
    ("profit",            "SUM(COALESCE(o.OrderTotal,0) - COALESCE(o.OrderSubtotalExclTax,0))"),
    ("line_item_revenue", "SUM(oi.PriceExclTax)"),
    ("tax_amount",        "SUM(o.OrderTax)"),
    ("line_item_cost",    "SUM(oi.OriginalProductCost)"),
    ("line_item_discount","SUM(oi.DiscountAmountExclTax)"),
    ("shipment_count",    "COUNT(DISTINCT sh.Id)"),
]

cols = 3
col_w = Inches(4.25)
for i, (mid, sql) in enumerate(metrics):
    row = i // cols
    col = i % cols
    x = Inches(0.3) + col * (col_w + Inches(0.1))
    y = Inches(1.45) + row * Inches(1.03)
    add_rect(s, x, y, col_w, Inches(0.98), fill=LIGHTBG, border=TEAL,
             line_w=Pt(0.75))
    add_rect(s, x, y, Inches(0.16), Inches(0.98), fill=TEAL)
    txb(s, mid,  x + Inches(0.22), y + Inches(0.04),
        col_w - Inches(0.28), Inches(0.38), size=13, bold=True, color=TEAL)
    txb(s, sql,  x + Inches(0.22), y + Inches(0.44),
        col_w - Inches(0.28), Inches(0.48), size=9.5, color=MID, italic=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — 34 DIMENSIONS & BFS JOINS
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "Semantic Layer — 34 Dimensions & BFS Join Resolution")
footer(s)

# Left: dimension categories
txb(s, "34 Dimensions by Category",
    Inches(0.35), Inches(1.45), Inches(6.1), Inches(0.4),
    size=17, bold=True, color=NAVY)

dim_categories = [
    (TEAL, "Products (8)",       "product_name, product_sku,\nproduct_price, product_cost,\nproduct_stock, product_rating,\nproduct_published, product_created_date"),
    (NAVY, "Orders — Time (5)",  "order_date, order_month,\norder_year, order_id, order_number"),
    (TEAL, "Orders — Status (9)","order_status, payment_status,\nshipping_status, payment_method,\ncurrency_code, shipping_method,\nOrderStatusId, PaymentStatusId, ShippingStatusId"),
    (NAVY, "Customers (4)",      "customer_name, customer_email,\ncustomer_active,\ncustomer_registration_date"),
    (TEAL, "Taxonomy (2)",       "category_name,\nmanufacturer_name"),
    (NAVY, "Shipments (3)",      "tracking_number,\nshipped_date, delivery_date"),
    (TEAL, "Geography (2)",      "country_name,\nbilling_city"),
    (NAVY, "Store (1)",          "store_name"),
]
for i, (clr, cat, dims) in enumerate(dim_categories):
    row, col = divmod(i, 2)
    x = Inches(0.35) + col * Inches(3.1)
    y = Inches(1.95) + row * Inches(1.18)
    add_rect(s, x, y, Inches(2.95), Inches(1.08), fill=WHITE, border=clr, line_w=Pt(1))
    add_rect(s, x, y, Inches(2.95), Inches(0.36), fill=clr)
    txb(s, cat, x + Inches(0.08), y + Inches(0.04),
        Inches(2.8), Inches(0.3), size=12, bold=True, color=WHITE)
    txb(s, dims, x + Inches(0.08), y + Inches(0.42),
        Inches(2.8), Inches(0.6), size=10, color=DARK)

# Right: BFS illustration
add_rect(s, Inches(6.55), Inches(1.45), Inches(0.06), Inches(5.35), fill=LIGHTGRAY)

txb(s, "BFS Join Resolution",
    Inches(6.8), Inches(1.45), Inches(6.2), Inches(0.4),
    size=17, bold=True, color=NAVY)
txb(s, 'Example: "revenue by category"',
    Inches(6.8), Inches(1.9), Inches(6.2), Inches(0.35),
    size=13.5, italic=True, color=TEAL)

# BFS path boxes
bfs_steps = [
    (TEAL,  "ORDER",                  "binding_table for metric: revenue"),
    (NAVY,  "→  ORDER_ITEM",          "JOIN Order_Item oi ON o.Id = oi.OrderId"),
    (TEAL,  "→  PRODUCT",             "JOIN Product p ON oi.ProductId = p.Id"),
    (NAVY,  "→  PRODUCT_CAT_MAPPING", "JOIN Product_Category_Mapping pcm ON p.Id = pcm.ProductId"),
    (TEAL,  "→  CATEGORY",            "JOIN Category c ON pcm.CategoryId = c.Id"),
]
for i, (clr, node, clause) in enumerate(bfs_steps):
    y = Inches(2.4) + i * Inches(0.82)
    add_rect(s, Inches(6.8), y, Inches(6.1), Inches(0.72), fill=LIGHTBG, border=clr)
    add_rect(s, Inches(6.8), y, Inches(0.14), Inches(0.72), fill=clr)
    txb(s, node,   Inches(7.0), y + Inches(0.04), Inches(3.0), Inches(0.32),
        size=13, bold=True, color=clr)
    txb(s, clause, Inches(7.0), y + Inches(0.38), Inches(5.7), Inches(0.28),
        size=10, color=MID, italic=True)

txb(s, "Only 11 pre-approved ON clauses exist — BFS finds the shortest path.",
    Inches(6.8), Inches(6.55), Inches(6.1), Inches(0.38),
    size=12, bold=True, color=GREENOK)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — SCOPE JUSTIFICATION (12 of 126 tables)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "Why Only 12 of 126 Tables? — Scope Justification",
           "The 114 hidden tables are not analytics tables — their exclusion is a security feature, not a gap")
footer(s)

# Left: the two groups
card(s, Inches(0.35), Inches(1.45), Inches(6.1), Inches(5.35), fill=WHITE, border=GREENOK)
add_rect(s, Inches(0.35), Inches(1.45), Inches(6.1), Inches(0.42), fill=GREENOK)
txb(s, "✅  12 Exposed — Analytics Domain",
    Inches(0.45), Inches(1.5), Inches(5.9), Inches(0.35),
    size=14, bold=True, color=WHITE)

exposed = [
    ("Orders",    "The core business object — revenue,\nstatus, payment, shipping"),
    ("Products",  "What was sold — name, SKU, price,\ncategory, manufacturer"),
    ("Customers", "Who bought — email, registration,\nlocation"),
    ("Geography", "Address → Country for billing region"),
    ("Shipments", "Fulfilment — tracking, ship/delivery date"),
    ("Store",     "Multi-store segmentation"),
]
for i, (label, desc) in enumerate(exposed):
    y = Inches(2.05) + i * Inches(0.77)
    add_rect(s, Inches(0.48), y, Inches(1.0), Inches(0.32), fill=GREENOK)
    txb(s, label, Inches(0.5), y + Inches(0.03), Inches(0.96), Inches(0.28),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, desc,  Inches(1.55), y + Inches(0.01), Inches(4.7), Inches(0.65),
        size=11, color=DARK)

add_rect(s, Inches(6.6), Inches(1.45), Inches(0.06), Inches(5.35), fill=LIGHTGRAY)

card(s, Inches(6.75), Inches(1.45), Inches(6.2), Inches(5.35), fill=WHITE, border=REDNG)
add_rect(s, Inches(6.75), Inches(1.45), Inches(6.2), Inches(0.42), fill=REDNG)
txb(s, "🔒  114 Hidden — Not Analytics Tables",
    Inches(6.85), Inches(1.5), Inches(6.0), Inches(0.35),
    size=14, bold=True, color=WHITE)

hidden = [
    ("System",    "Log, ScheduleTask, Setting, ActivityLog,\nGenericAttribute"),
    ("CMS",       "Topic, NewsItem, BlogPost,\nPoll, PollAnswer"),
    ("Config",    "Language, Currency, LocalizedProperty,\nMeasureUnit"),
    ("Security",  "PermissionRecord, AclRecord,\nCustomerPassword"),
    ("Pre-order", "ShoppingCartItem, WishlistItem"),
    ("Promos",    "GiftCard, Discount, RewardPoints,\nDiscountUsageHistory"),
]
for i, (label, desc) in enumerate(hidden):
    y = Inches(2.05) + i * Inches(0.77)
    add_rect(s, Inches(6.88), y, Inches(1.0), Inches(0.32), fill=REDNG)
    txb(s, label, Inches(6.9), y + Inches(0.03), Inches(0.96), Inches(0.28),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, desc,  Inches(7.95), y + Inches(0.01), Inches(4.8), Inches(0.65),
        size=11, color=DARK)

txb(s, "Key point: 114 unexposed tables include security-sensitive data (passwords, ACLs) "
       "that must NOT be exposed to a self-service analytics tool. Exclusion = access control.",
    Inches(0.35), Inches(7.0), Inches(12.6), Inches(0.35),
    size=11.5, color=NAVY, italic=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — VOCABULARY INJECTION
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "Vocabulary Injection — How the LLM Understands Business Terms",
           'Why SYNONYMS = {}  and why that is a feature, not a bug')
footer(s)

# Left side: the mechanism
txb(s, "How it works",
    Inches(0.35), Inches(1.45), Inches(5.8), Inches(0.38),
    size=18, bold=True, color=NAVY)

steps = [
    ("1", TEAL,  "Approved vocabulary embedded in system prompt",
                 "All 15 metric labels + 34 dimension labels (~1,100 tokens)\nare injected into the LLM's system prompt at every request."),
    ("2", GOLD,  "LLM performs fuzzy matching at inference time",
                 '"Show me sales by product line" → LLM maps:\n  sales  →  revenue\n  product line  →  category_name'),
    ("3", GREENOK,"Output is a validated IntentObject — no SQL",
                 '{"metric": "revenue", "dimension": "category_name",\n "filter": null, "pattern": "ranking"}'),
]
for i, (num, clr, title, desc) in enumerate(steps):
    y = Inches(1.92) + i * Inches(1.6)
    add_rect(s, Inches(0.35), y, Inches(0.6), Inches(0.6), fill=clr)
    txb(s, num, Inches(0.35), y + Inches(0.04), Inches(0.6), Inches(0.52),
        size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, title, Inches(1.05), y + Inches(0.04), Inches(5.1), Inches(0.4),
        size=14.5, bold=True, color=clr)
    txb(s, desc,  Inches(1.05), y + Inches(0.48), Inches(5.1), Inches(1.0),
        size=12.5, color=DARK, italic=False)

# Right side: SYNONYMS = {} explanation
add_rect(s, Inches(6.55), Inches(1.4), Inches(0.06), Inches(5.4), fill=LIGHTGRAY)

add_rect(s, Inches(6.8), Inches(1.45), Inches(6.1), Inches(1.55), fill=NAVY)
txb(s, "SYNONYMS = {}", Inches(6.9), Inches(1.52), Inches(5.9), Inches(0.5),
    size=28, bold=True, color=GOLD, font="Courier New")
txb(s, "Intentionally empty — and this is a design strength.",
    Inches(6.9), Inches(2.0), Inches(5.9), Inches(0.4),
    size=13, color=TEALLIGHT)
txb(s, "Why this is better than a synonym dictionary:",
    Inches(6.8), Inches(3.12), Inches(6.1), Inches(0.38),
    size=14.5, bold=True, color=NAVY)

adv = [
    "No maintenance burden — new business terms need no code change",
    "LLM handles linguistic variation naturally (plurals, abbreviations, jargon)",
    "Synonym dictionary grows stale; vocabulary injection stays current",
    "Reduces the semantic layer to pure SQL logic — cleaner separation of concerns",
]
bullet_box(s, [(a, 13, DARK, False) for a in adv],
           Inches(6.8), Inches(3.55), Inches(6.1), Inches(2.6))

# business logic box
add_rect(s, Inches(6.8), Inches(6.25), Inches(6.1), Inches(0.58), fill=LIGHTBLUE, border=TEAL)
txb(s, 'Business Logic Mapping:  "abandoned"  →  OrderStatusId = 40',
    Inches(6.9), Inches(6.33), Inches(5.9), Inches(0.42),
    size=13, bold=True, color=NAVY, font="Courier New")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — QUERY WALKTHROUGH
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "End-to-End: How a Query Flows Through AEGIS",
           '"Show me revenue by product category for this month"')
footer(s)

flow = [
    (TEAL, "Stage 1 — Intent Extraction (LLM)",
     'User: "Show me revenue by product category for this month"\n'
     'LLM output →  {"metric":"revenue", "dimension":"category_name",\n'
     '               "filter":{"order_month":"current"}, "pattern":"ranking"}'),
    (GOLD, "Stage 2 — Coverage Validation",
     '"revenue" ∈ METRICS ✓    "category_name" ∈ DIMENSIONS ✓\n'
     'Pattern "ranking" ∈ ANALYTICS_PRIMITIVES ✓    Query accepted.'),
    (NAVY, "Stage 3 — Semantic Mapping",
     'Resolves "revenue" → Metric object with sql_expr, binding_table="Order"\n'
     'Resolves "category_name" → Dimension with required_joins=[...]\n'
     'Expands filter: order_month="current" → DATE_FORMAT(...) = current month'),
    (TEAL, "Stage 4 — Permission Rewriting",
     'Appends row-level security:  AND o.StoreId IN (1,2)  [if user is store-scoped]\n'
     'Admin users: no additional WHERE predicates appended.'),
    (GREENOK, "Stage 5 — SQL Compilation (BFS + Template)",
     'BFS: Order → OrderItem → Product → Product_Category_Mapping → Category\n'
     'SELECT c.Name, SUM(o.OrderTotal-o.RefundedAmount) AS revenue\n'
     'FROM Order o JOIN ... WHERE ... GROUP BY c.Name ORDER BY revenue DESC'),
    (NAVY, "Stages 6+7 — Visualization & Widget Persistence",
     'Pattern="ranking" → bar_chart selected automatically\n'
     'SHA-256 hash: widget already exists → returns cached widget ID'),
]

for i, (clr, title, detail) in enumerate(flow):
    y = Inches(1.42) + i * Inches(0.99)
    add_rect(s, Inches(0.35), y, Inches(12.6), Inches(0.9), fill=WHITE,
             border=clr, line_w=Pt(1))
    add_rect(s, Inches(0.35), y, Inches(0.18), Inches(0.9), fill=clr)
    txb(s, title,  Inches(0.6), y + Inches(0.04), Inches(4.0), Inches(0.36),
        size=12.5, bold=True, color=clr)
    txb(s, detail, Inches(4.7), y + Inches(0.02), Inches(8.1), Inches(0.84),
        size=11, color=DARK, italic=False, font="Courier New")

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — SQL SAFETY GUARANTEE
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "SQL Safety — Formal Guarantee + Two-Layer Defense")
footer(s)

# Formal claim box
add_rect(s, Inches(0.35), Inches(1.45), Inches(12.6), Inches(1.4), fill=NAVY)
txb(s, "Formal Safety Claim",
    Inches(0.5), Inches(1.52), Inches(12.3), Inches(0.38),
    size=14, bold=True, color=GOLD)
txb(s, "Given that the SQL compiler only accepts validated IntentObject inputs and generates SQL exclusively "
       "from pre-defined templates, the set of possible SQL outputs is finite and enumerable. "
       "SQL injection requires generating SQL outside this set — which is architecturally impossible.",
    Inches(0.5), Inches(1.92), Inches(12.3), Inches(0.8),
    size=13.5, color=WHITE, italic=True)

# Two layers
txb(s, "Two-Layer SQL Safety",
    Inches(0.35), Inches(3.0), Inches(12.6), Inches(0.38),
    size=17, bold=True, color=NAVY)

layer1_pts = [
    "All SQL assembled from named template fragments",
    "User-controlled values bound as SQL parameters — never interpolated",
    "Column and table names are compile-time constants from the semantic layer",
    "LLM output (IntentObject JSON) never concatenated into SQL string",
    "Injection requires escaping the template system — structurally impossible",
]
card(s, Inches(0.35), Inches(3.45), Inches(5.95), Inches(3.25), fill=RGBColor(0xF0,0xF8,0xF0), border=GREENOK)
add_rect(s, Inches(0.35), Inches(3.45), Inches(5.95), Inches(0.42), fill=GREENOK)
txb(s, "Layer 1 — Parameterized Template Compilation  (Structural Prevention)",
    Inches(0.45), Inches(3.5), Inches(5.75), Inches(0.35),
    size=11.5, bold=True, color=WHITE)
bullet_box(s, [(p, 12.5, DARK, False) for p in layer1_pts],
           Inches(0.45), Inches(3.92), Inches(5.75), Inches(2.65))

layer2_pts = [
    "Post-compilation scanner checks 16 forbidden patterns",
    "DROP, DELETE, INSERT, UPDATE, TRUNCATE, ALTER, CREATE",
    "UNION, comment sequences (--, /*), hex literals, stacked queries (;)",
    "Catches any hypothetical compiler bug before execution",
    "Defense-in-depth — belt AND suspenders",
]
card(s, Inches(6.7), Inches(3.45), Inches(6.2), Inches(3.25), fill=RGBColor(0xF0,0xF8,0xF0), border=TEAL)
add_rect(s, Inches(6.7), Inches(3.45), Inches(6.2), Inches(0.42), fill=TEAL)
txb(s, "Layer 2 — Post-Compilation Safety Scanner  (Defense-in-Depth)",
    Inches(6.8), Inches(3.5), Inches(6.0), Inches(0.35),
    size=11.5, bold=True, color=WHITE)
bullet_box(s, [(p, 12.5, DARK, False) for p in layer2_pts],
           Inches(6.8), Inches(3.92), Inches(6.0), Inches(2.65))

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — EVALUATION RESULTS
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "Evaluation Results — 100-Query Benchmark",
           "nopCommerce e-commerce domain  |  5 systems compared across 3 safety & accuracy metrics")
footer(s)

# Table
headers = ["System", "Unsafe SQL", "Execution Validity", "Coverage"]
rows = [
    ("B1: Direct GPT-4",        "5.0%",   "99.0%", "99.0%",  REDNG,   WHITE),
    ("B2: Decomposed LLM",        "3.0%",   "97.0%", "97.0%",  REDNG,   WHITE),
    ("B3: Template-only (no LLM)","1.0%",  "66.0%", "55.0%",  REDNG,   WHITE),
    ("B4: AEGIS ablated",        "0.0%",   "88.7%", "91.0%",  GOLD,    DARK),
    ("AEGIS (This Thesis)",      "0.0%",   "100.0%","100.0%", GREENOK, WHITE),
]

th = Inches(0.42)  # table header height
tr = Inches(0.72)  # table row height
tx = Inches(0.5)
ty = Inches(1.55)
tw = [Inches(3.6), Inches(2.2), Inches(2.7), Inches(2.7)]
cols_x = [tx]
for w in tw[:-1]:
    cols_x.append(cols_x[-1] + w)

# header row
for j, (hdr, cx, cw) in enumerate(zip(headers, cols_x, tw)):
    add_rect(s, cx, ty, cw, th, fill=NAVY)
    txb(s, hdr, cx + Inches(0.08), ty + Inches(0.06), cw - Inches(0.16), th - Inches(0.1),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, (system, unsafe, validity, coverage, badge_clr, badge_txt_clr) in enumerate(rows):
    ry = ty + th + i * tr
    row_fill = RGBColor(0xF0, 0xFF, 0xF4) if i == 4 else (LIGHTBG if i % 2 == 0 else WHITE)
    border_clr = GREENOK if i == 4 else LIGHTGRAY
    cells = [system, unsafe, validity, coverage]
    for j, (val, cx, cw) in enumerate(zip(cells, cols_x, tw)):
        add_rect(s, cx, ry, cw, tr, fill=row_fill, border=border_clr, line_w=Pt(0.75))
        cell_color = DARK
        cell_bold = (i == 4)
        if j == 1 and val != "0.0%":   # unsafe SQL — bad
            cell_color = REDNG
        elif j == 1 and val == "0.0%": # unsafe SQL — good
            cell_color = GREENOK
        elif j in (2, 3) and val == "100.0%":
            cell_color = GREENOK
        elif j in (2, 3) and float(val.replace('%','')) < 90:
            cell_color = REDNG
        txb(s, val, cx + Inches(0.08), ry + Inches(0.14), cw - Inches(0.16), tr - Inches(0.2),
            size=15 if i == 4 else 14, bold=cell_bold, color=cell_color,
            align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

# AEGIS highlight arrow
add_rect(s, Inches(0.5), ty + th + 4 * tr, Inches(11.2), tr,
         fill=None, border=GREENOK, line_w=Pt(2))

# insight boxes
insights = [
    (GREENOK, "0% Unsafe SQL", "Not a statistical finding —\na structural guarantee from\nthe architecture."),
    (TEAL,    "100% Validity",  "Every query compiles and\nexecutes without error on the\nfirst attempt."),
    (NAVY,    "100% Coverage",  "Every benchmark question is\nwithin the semantic layer\nvocabulary."),
]
for i, (clr, title, desc) in enumerate(insights):
    x = Inches(0.5) + i * Inches(4.0)
    y = Inches(6.1)
    add_rect(s, x, y, Inches(3.7), Inches(1.12), fill=WHITE, border=clr, line_w=Pt(1.5))
    add_rect(s, x, y, Inches(3.7), Inches(0.38), fill=clr)
    txb(s, title, x + Inches(0.1), y + Inches(0.05), Inches(3.5), Inches(0.32),
        size=13.5, bold=True, color=WHITE)
    txb(s, desc,  x + Inches(0.1), y + Inches(0.44), Inches(3.5), Inches(0.62),
        size=11.5, color=DARK)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — ABLATION STUDY
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "Ablation Study — Contribution of Each Component",
           "Removing each component one at a time to measure its individual contribution")
footer(s)

ablations = [
    ("Full AEGIS",                  100.0, 100.0, GREENOK),
    ("− Permission Rewriting",       99.0,  99.0, TEAL),
    ("− Post-compilation Scanner",   99.0,  99.0, TEAL),
    ("− BFS Join Resolution",        88.7,  91.0, GOLD),
    ("− Business Logic Mappings",    95.0,  93.0, GOLD),
    ("− Vocabulary Injection",       64.7,  65.0, REDNG),
    ("LLM Only (No Semantic Layer)", 12.0,   9.0, REDNG),
]

# bar chart (manual rectangles)
bar_area_x = Inches(5.2)
bar_area_y = Inches(1.5)
bar_area_w = Inches(7.7)
bar_area_h = Inches(5.3)
max_val = 105.0

add_rect(s, bar_area_x, bar_area_y, bar_area_w, bar_area_h, fill=LIGHTBG, border=LIGHTGRAY)

bar_count = len(ablations)
bar_h_each = bar_area_h / (bar_count + 1)
bar_max_w = bar_area_w - Inches(0.3)

for i, (label, validity, coverage, clr) in enumerate(ablations):
    by = bar_area_y + bar_h_each * (i + 0.4)
    # validity bar
    bw_v = bar_max_w * (validity / max_val)
    add_rect(s, bar_area_x + Inches(0.1), by,
             bw_v, bar_h_each * 0.38, fill=clr)
    # coverage bar (slightly offset)
    bw_c = bar_max_w * (coverage / max_val)
    lighter = RGBColor(
        min(255, clr[0] + 40), min(255, clr[1] + 40), min(255, clr[2] + 40))
    add_rect(s, bar_area_x + Inches(0.1), by + bar_h_each * 0.4,
             bw_c, bar_h_each * 0.35, fill=lighter)
    # value labels
    txb(s, f"{validity:.0f}%", bar_area_x + Inches(0.12) + bw_v, by - Inches(0.02),
        Inches(0.7), bar_h_each * 0.38, size=10, bold=(i==0), color=clr)

# row labels on left
for i, (label, validity, coverage, clr) in enumerate(ablations):
    by = bar_area_y + bar_h_each * (i + 0.4)
    txb(s, label, Inches(0.3), by + Inches(0.03),
        Inches(4.75), bar_h_each * 0.75, size=13,
        bold=(i==0), color=DARK if i > 0 else GREENOK)

# legend
add_rect(s, Inches(5.5), Inches(6.65), Inches(0.32), Inches(0.2), fill=TEAL)
txb(s, "Execution Validity", Inches(5.85), Inches(6.62), Inches(2.0), Inches(0.28),
    size=11, color=DARK)
add_rect(s, Inches(8.0), Inches(6.65), Inches(0.32), Inches(0.2),
         fill=RGBColor(0x70, 0xD0, 0xFF))
txb(s, "Coverage", Inches(8.35), Inches(6.62), Inches(1.5), Inches(0.28),
    size=11, color=DARK)

# key finding callout
add_rect(s, Inches(0.3), Inches(6.7), Inches(4.85), Inches(0.52), fill=NAVY)
txb(s, "Vocabulary injection contributes −35.3pp validity — the single most impactful component.",
    Inches(0.4), Inches(6.76), Inches(4.7), Inches(0.42),
    size=11.5, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 15 — GENERALIZABILITY
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "Generalizability — WooCommerce Schema Evaluation (RQ4)",
           "Does AEGIS work beyond the nopCommerce schema it was built on?")
footer(s)

# Left: what was done
card(s, Inches(0.35), Inches(1.45), Inches(5.85), Inches(5.35),
     fill=WHITE, border=TEAL)
add_rect(s, Inches(0.35), Inches(1.45), Inches(5.85), Inches(0.45), fill=TEAL)
txb(s, "What Was Evaluated",
    Inches(0.45), Inches(1.5), Inches(5.65), Inches(0.38),
    size=15, bold=True, color=WHITE)
steps_gen = [
    ("Schema",      "WooCommerce database — completely different tables,\ncolumn names, and relationships from nopCommerce."),
    ("Method",      "Built a new semantic layer (12 metrics, 28 dims,\n9 join paths) from scratch. Zero pipeline changes."),
    ("Benchmark",   "Same 100-query evaluation dataset, re-run on\nthe WooCommerce semantic layer."),
    ("Constraint",  "Single developer, measured real elapsed time\nto build the new semantic layer end-to-end."),
]
for i, (label, desc) in enumerate(steps_gen):
    y = Inches(2.05) + i * Inches(1.12)
    add_rect(s, Inches(0.48), y, Inches(1.1), Inches(0.38), fill=NAVY)
    txb(s, label, Inches(0.5), y + Inches(0.04), Inches(1.05), Inches(0.3),
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, desc, Inches(1.65), y + Inches(0.02), Inches(4.4), Inches(1.0),
        size=12.5, color=DARK)

# Right: results
card(s, Inches(7.1), Inches(1.45), Inches(5.85), Inches(5.35),
     fill=WHITE, border=GREENOK)
add_rect(s, Inches(7.1), Inches(1.45), Inches(5.85), Inches(0.45), fill=GREENOK)
txb(s, "Results",
    Inches(7.2), Inches(1.5), Inches(5.65), Inches(0.38),
    size=15, bold=True, color=WHITE)

result_metrics = [
    ("98.0%",  "Intent Accuracy",      GREENOK),
    ("0%",     "Unsafe SQL Queries",   GREENOK),
    ("14 hrs", "Semantic Layer Build", TEAL),
    ("0",      "Pipeline Code Changes",TEAL),
]
for i, (val, label, clr) in enumerate(result_metrics):
    x = Inches(7.2) + (i % 2) * Inches(2.95)
    y = Inches(2.1) + (i // 2) * Inches(1.85)
    add_rect(s, x, y, Inches(2.75), Inches(1.65), fill=LIGHTBG, border=clr)
    txb(s, val, x + Inches(0.1), y + Inches(0.18), Inches(2.55), Inches(0.75),
        size=36, bold=True, color=clr, align=PP_ALIGN.CENTER)
    txb(s, label, x + Inches(0.1), y + Inches(1.1), Inches(2.55), Inches(0.45),
        size=12.5, color=DARK, align=PP_ALIGN.CENTER)

txb(s, "Conclusion: AEGIS is a reusable architecture, not a one-schema prototype.\n"
       "Any e-commerce schema can be onboarded by writing a new semantic layer —\n"
       "no retraining, no pipeline changes, no LLM reconfiguration.",
    Inches(7.2), Inches(5.7), Inches(5.65), Inches(0.9),
    size=12.5, color=DARK, italic=True)

# divider
add_rect(s, Inches(6.55), Inches(1.45), Inches(0.06), Inches(5.35), fill=LIGHTGRAY)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 16 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "Conclusion & Future Work")
footer(s)

# contributions
txb(s, "Key Contributions",
    Inches(0.35), Inches(1.45), Inches(12.6), Inches(0.38),
    size=18, bold=True, color=NAVY)

contribs = [
    (TEAL,    "Safety-by-Design Architecture",
              "First NL2SQL system that makes SQL injection structurally impossible by separating the AI layer from the SQL generation layer."),
    (NAVY,    "Semantic Layer with Vocabulary Injection",
              "Closed vocabulary of 15 metrics, 34 dimensions, 11 join paths — with dynamic vocabulary injection replacing static synonym dictionaries."),
    (GOLD,    "Persistent Widget Engine",
              "Turns one-off query results into refreshable dashboard widgets with SHA-256 deduplication and scheduled refresh (61% recurring)."),
    (GREENOK, "Formal Safety Proof + Empirical Validation",
              "Mathematical proof that injection is impossible, validated empirically with 0% unsafe SQL across 100 benchmark queries."),
]
for i, (clr, title, desc) in enumerate(contribs):
    row, col = divmod(i, 2)
    x = Inches(0.35) + col * Inches(6.5)
    y = Inches(1.95) + row * Inches(1.35)
    add_rect(s, x, y, Inches(6.15), Inches(1.22), fill=WHITE, border=clr, line_w=Pt(1.5))
    add_rect(s, x, y, Inches(0.18), Inches(1.22), fill=clr)
    txb(s, title, x + Inches(0.28), y + Inches(0.06), Inches(5.7), Inches(0.38),
        size=14.5, bold=True, color=clr)
    txb(s, desc,  x + Inches(0.28), y + Inches(0.48), Inches(5.7), Inches(0.68),
        size=12.5, color=DARK)

# future work
txb(s, "Future Work",
    Inches(0.35), Inches(4.75), Inches(12.6), Inches(0.38),
    size=18, bold=True, color=NAVY)
future = [
    "Extend semantic layer to transactional databases (banking, healthcare, logistics)",
    "Real-time semantic layer editor — add new metrics/dimensions without code deployment",
    "Multi-turn conversational analytics — follow-up questions building on previous widgets",
    "Fine-tune smaller LLM on IntentObject extraction to reduce API latency from 1.85s",
]
bullet_box(s, [(f, 14, DARK, False) for f in future],
           Inches(0.35), Inches(5.18), Inches(12.6), Inches(1.7))

# ═══════════════════════════════════════════════════════════════════
# SLIDE 17 — WOOCOMMERCE DEPLOYMENT GUIDE
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s)
header_bar(s, "How a WooCommerce Store Owner Deploys AEGIS",
           "Only the semantic layer changes — the entire pipeline is reused without modification")
footer(s)

# Key principle banner
add_rect(s, Inches(0.35), Inches(1.45), Inches(12.6), Inches(0.55), fill=NAVY)
txb(s, "Core rule: write a new semantic_layer.py for your schema. Touch nothing else.",
    Inches(0.5), Inches(1.52), Inches(12.3), Inches(0.42),
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# Steps — left column
steps_left = [
    ("1", TEAL,  "Prerequisites  (30 min)",
     "Python 3.10+, pip, Git, MySQL read access\nGroq API key (free tier sufficient)\ngit clone + pip install -r requirements.txt"),
    ("2", NAVY,  "Schema Analysis  (2–3 hrs)",
     "Map WooCommerce tables to AEGIS concepts:\n  KPIs the business tracks  →  METRICS\n  Slice-by axes  →  DIMENSIONS\n  Table relationships  →  JOIN_GRAPH"),
    ("3", GOLD,  "Build the Semantic Layer  (8–10 hrs)",
     "WooCommerce result: 12 metrics, 28 dimensions,\n9 join paths, 18 tables.\nOnly file that changes: semantic_layer.py"),
]
for i, (num, clr, title, desc) in enumerate(steps_left):
    y = Inches(2.12) + i * Inches(1.52)
    add_rect(s, Inches(0.35), y, Inches(0.58), Inches(0.58), fill=clr)
    txb(s, num, Inches(0.35), y + Inches(0.04), Inches(0.58), Inches(0.5),
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, title, Inches(1.05), y + Inches(0.04), Inches(5.5), Inches(0.38),
        size=13.5, bold=True, color=clr)
    txb(s, desc,  Inches(1.05), y + Inches(0.45), Inches(5.5), Inches(1.0),
        size=12, color=DARK)

# Steps — right column
steps_right = [
    ("4", TEAL,  "Configure Database  (30 min)",
     "Set DB_HOST, DB_NAME, DB_USER, DB_PASSWORD\nand LLM_BASE_URL + LLM_API_KEY in .env\n(any OpenAI-compatible endpoint)"),
    ("5", NAVY,  "Test  (1–2 hrs)",
     'python -m unittest discover -s tests\npython run_demo_cli.py\n"Show me revenue by category this month"'),
    ("6", GREENOK,"Deploy  (30 min)",
     "docker-compose up --build\nDashboard at http://localhost:8765\nOr: python run_demo_server.py"),
]
for i, (num, clr, title, desc) in enumerate(steps_right):
    y = Inches(2.12) + i * Inches(1.52)
    add_rect(s, Inches(7.1), y, Inches(0.58), Inches(0.58), fill=clr)
    txb(s, num, Inches(7.1), y + Inches(0.04), Inches(0.58), Inches(0.5),
        size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txb(s, title, Inches(7.8), y + Inches(0.04), Inches(5.3), Inches(0.38),
        size=13.5, bold=True, color=clr)
    txb(s, desc,  Inches(7.8), y + Inches(0.45), Inches(5.3), Inches(1.0),
        size=12, color=DARK, font="Courier New")

# Divider
add_rect(s, Inches(6.6), Inches(2.05), Inches(0.06), Inches(4.65), fill=LIGHTGRAY)

# Effort table at bottom
add_rect(s, Inches(0.35), Inches(6.72), Inches(12.6), Inches(0.5), fill=LIGHTBG, border=TEAL)
effort_cols = [
    ("Schema analysis", "2–3 hrs"),
    ("Metrics (SQL expressions)", "2–3 hrs"),
    ("Dimensions + join paths", "4–5 hrs"),
    ("Test & iterate", "2 hrs"),
    ("Total", "~14 hrs"),
]
col_w = Inches(2.52)
for i, (label, time) in enumerate(effort_cols):
    x = Inches(0.35) + i * col_w
    bg = NAVY if label == "Total" else LIGHTBG
    add_rect(s, x, Inches(6.72), col_w, Inches(0.5), fill=bg)
    txb(s, f"{label}:  {time}", x + Inches(0.08), Inches(6.79),
        col_w - Inches(0.16), Inches(0.36),
        size=11.5, bold=(label=="Total"),
        color=GOLD if label=="Total" else DARK,
        align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 18 — THANK YOU
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, NAVY)
add_rect(s, 0, 0, Inches(0.55), SLIDE_H, fill=GOLD)
add_rect(s, Inches(0.55), 0, Inches(0.08), SLIDE_H, fill=TEAL)

txb(s, "Thank You",
    Inches(1.1), Inches(1.4), Inches(11.4), Inches(1.2),
    size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Calibri Light")

add_rect(s, Inches(2.5), Inches(2.7), Inches(8.3), Inches(0.055), fill=GOLD)

txb(s, "Questions & Discussion",
    Inches(1.1), Inches(2.85), Inches(11.4), Inches(0.6),
    size=26, color=TEALLIGHT, align=PP_ALIGN.CENTER, font="Calibri Light")

summary_pts = [
    "0% SQL injection — structural guarantee, not a detection heuristic",
    "100% execution validity — every query correct on first attempt",
    "98% accuracy on unseen schema (WooCommerce) — generalizable architecture",
    "14 person-hours to port to a new schema — practical for real deployment",
]
for i, pt in enumerate(summary_pts):
    y = Inches(3.65) + i * Inches(0.68)
    add_rect(s, Inches(2.0), y + Inches(0.14), Inches(0.22), Inches(0.22), fill=GOLD)
    txb(s, pt, Inches(2.35), y, Inches(9.0), Inches(0.6),
        size=15.5, color=WHITE)

add_rect(s, Inches(2.5), Inches(6.55), Inches(8.3), Inches(0.035), fill=GOLD)
txb(s, "Md. Riaz  |  Pundra University of Science and Technology  |  BSc CSE  |  2026",
    Inches(1.1), Inches(6.65), Inches(11.4), Inches(0.45),
    size=13, color=RGBColor(0x8A,0xA8,0xC8), align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────
import os
out_path = os.path.join(os.path.dirname(__file__), "..", "AEGIS_Thesis_Defense.pptx")
out_path = os.path.normpath(out_path)
prs.save(out_path)
print(f"Saved -> {out_path}")
print(f"Slides: {len(prs.slides)}")
