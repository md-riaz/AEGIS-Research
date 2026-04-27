import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # helper for adding slides
    def add_slide(title_text, content_text=None, image_path=None, bullet_points=None):
        slide_layout = prs.slide_layouts[1] if content_text or bullet_points else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        if bullet_points:
            tf = slide.placeholders[1].text_frame
            tf.text = bullet_points[0]
            for point in bullet_points[1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
        elif content_text:
            slide.placeholders[1].text = content_text

        if image_path and os.path.exists(image_path):
            # Adjust position based on whether there is text
            left = Inches(5.5) if bullet_points or content_text else Inches(2)
            top = Inches(2)
            height = Inches(4)
            slide.shapes.add_picture(image_path, left, top, height=height)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Natural Language to Dashboard"
    subtitle = slide.placeholders[1]
    subtitle.text = "A Safe AI-Assisted Reporting and Widget Generation System\n\nMd. Riaz\nDepartment of CSE, Pundra University"

    # Slide 2: Introduction - The Problem
    add_slide("The Problem: Reporting is Slow", bullet_points=[
        "Organizations have data, but non-technical staff can't access it easily.",
        "Average wait time for a new report: 3.2 days.",
        "61% of requests are just variations of old questions.",
        "Traditional AI (Text-to-SQL) is risky: it can make mistakes or expose data."
    ])

    # Slide 3: The Goal - SafeDash
    add_slide("The Goal: SafeDash", bullet_points=[
        "Turn plain English into persistent, reusable widgets.",
        "Focus on safety: The AI understands the question, but templates build the SQL.",
        "Dynamic and refreshable: Save once, use every day.",
        "No manually written synonyms needed."
    ])

    # Slide 4: System Architecture
    add_slide("System Architecture", 
              bullet_points=[
                  "1. User asks a question in plain English.",
                  "2. AI Parser extracts intent and business terms.",
                  "3. Semantic Layer maps terms to database logic.",
                  "4. Safe Compiler builds SQL from approved templates.",
                  "5. Widget Engine saves the result for reuse."
              ],
              image_path=r"C:\Users\mdriaz\.gemini\antigravity\brain\ba6797cc-7cf6-4fea-9f1e-eeda9b965eb5\safedash_architecture_diagram_1777290529737.png")

    # Slide 5: The Semantic Layer
    add_slide("The Heart: Semantic Layer", bullet_points=[
        "Defines Metrics: Revenue, Profit, Order Count.",
        "Defines Dimensions: Category, Time, Department.",
        "Join Paths: Pre-approved ways to connect tables.",
        "Security: Who is allowed to see what.",
        "Vocabulary Injection: AI sees these terms directly."
    ])

    # Slide 6: Dynamic Vocabulary Injection
    add_slide("Vocabulary Injection", bullet_points=[
        "Traditional way: Write long synonym lists (fragile).",
        "SafeDash way: Inject approved names + descriptions into the prompt.",
        "AI maps 'earnings' -> 'revenue' automatically.",
        "Zero maintenance: New metrics are instantly available.",
        "Results: 100% coverage with 0 manual synonyms."
    ])

    # Slide 7: Safe SQL Generation
    add_slide("Safe SQL Generation", bullet_points=[
        "AI never writes raw SQL strings.",
        "System uses Parameterized Templates (Ranking, Trend, KPI).",
        "AST Validation: Final check for forbidden commands (e.g., DELETE).",
        "structural impossibility of SQL injection.",
        "Reliable output every time."
    ])

    # Slide 8: Widget Persistence
    add_slide("Dashboard & Widgets", 
              bullet_points=[
                  "Widgets are first-class citizens.",
                  "Saved with: Analysis Plan, SQL, Chart Settings.",
                  "Reuse: System finds similar past widgets automatically.",
                  "Refreshable: Data updates without asking the AI again.",
                  "Premium UI: Glassmorphic and responsive."
              ],
              image_path=r"C:\Users\mdriaz\.gemini\antigravity\brain\ba6797cc-7cf6-4fea-9f1e-eeda9b965eb5\safedash_dashboard_preview_1777290576414.png")

    # Slide 9: Comparison
    add_slide("How SafeDash Compares", bullet_points=[
        "Most systems only do Text-to-SQL (Spider/BIRD).",
        "Some only do Visualization (nl4dv).",
        "SafeDash covers the full pipeline: NL -> SQL -> Vis -> Widget.",
        "First system to combine Semantic Layer + LLM + Persistence."
    ])

    # Slide 10: Evaluation Results
    add_slide("Evaluation: 100-Query Benchmark", 
              bullet_points=[
                  "Tested on real e-commerce data (nopCommerce).",
                  "100% Execution Validity.",
                  "100% Query Coverage.",
                  "0% Unsafe Queries (Baseline AI was 4.0% unsafe).",
                  "Zero manual synonyms required."
              ],
              image_path=r"C:\Users\mdriaz\.gemini\antigravity\brain\ba6797cc-7cf6-4fea-9f1e-eeda9b965eb5\safedash_performance_chart_1777290651053.png")

    # Slide 11: Future Work
    add_slide("Future Work", bullet_points=[
        "Support for multi-turn conversations.",
        "Automated semantic layer expansion from query logs.",
        "Support for cross-entity causal analysis.",
        "Mobile dashboard companion app."
    ])

    # Slide 12: Conclusion
    add_slide("Conclusion", bullet_points=[
        "SafeDash makes reporting safe, accurate, and easy.",
        "System design > Smarter AI for safety.",
        "Vocabulary injection removes the need for synonym lists.",
        "Widgets make analytics part of the daily workflow.",
        "Thank you! Any questions?"
    ])

    output_path = "SafeDash_Thesis_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
